"""Real PowerPoint and Excel file creation (python-pptx / openpyxl).

The files written here are opened again by the graders (and by the human
reviewer in actual PowerPoint/Excel), so this is genuine capability, not
a simulation.
"""
import os

from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .world import ToolError

# --- deck design -------------------------------------------------------------
# One small palette and a handful of measurements, applied consistently. The
# tool contract is unchanged ({"title", "bullets"}); only the rendering is.
INK = RGBColor(0x1F, 0x24, 0x30)      # near-black, for titles
BODY = RGBColor(0x3D, 0x45, 0x55)     # softer, for bullet text
MUTED = RGBColor(0x8A, 0x93, 0xA3)    # footers, slide numbers
ACCENT = RGBColor(0x2F, 0x6F, 0xEB)   # the one colour that does the work
FONT = "Calibri"                      # present on Windows and macOS Office

MARGIN = Inches(0.9)


def _style_run(para, size, color, bold=False, font=FONT):
    """Style the paragraph *and* every run in it.

    Paragraph-level font is only a default that runs inherit; anything reading
    run properties (our own web viewer, other pptx parsers) sees nothing. Being
    explicit at both levels makes the file say what it means.
    """
    for target in [para.font] + [r.font for r in para.runs]:
        target.size = Pt(size)
        target.color.rgb = color
        target.bold = bold
        target.name = font


def _bullet_size(n):
    """Shrink as a slide fills up, so a long list still fits the body box."""
    return 20 if n <= 5 else (17 if n <= 8 else 14)


def _resolve(files_dir, filename, ext):
    if not filename or not isinstance(filename, str):
        raise ToolError(f"'filename' must be a string ending in {ext}")
    name = os.path.basename(filename.strip())
    if not name.lower().endswith(ext):
        name += ext
    return os.path.join(files_dir, name), name


def create_presentation(files_dir, filename, slides):
    path, name = _resolve(files_dir, filename, ".pptx")
    if not isinstance(slides, list) or not slides:
        raise ToolError("'slides' must be a non-empty list of objects like "
                        '{"title": "...", "bullets": ["...", "..."]}')
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)  # 16:9
    W, H = prs.slide_width, prs.slide_height
    total = len(slides)

    for i, s in enumerate(slides):
        if not isinstance(s, dict) or "title" not in s:
            raise ToolError(f"slide {i + 1} must be an object with a 'title' key "
                            f"(and optional 'bullets' list), got {s!r}")
        bullets = s.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        if not isinstance(bullets, list):
            raise ToolError(f"slide {i + 1}: 'bullets' must be a list of strings")
        bullets = [str(b) for b in bullets]
        title_text = str(s["title"])

        if i == 0 and not bullets:
            # Cover. Layout 0 keeps a real title placeholder, so graders and the
            # web viewer can still find it via slide.shapes.title.
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT
            bar.line.fill.background()

            title = slide.shapes.title
            title.left, title.top = MARGIN, Inches(2.6)
            title.width, title.height = W - MARGIN * 2, Inches(1.5)
            title.text_frame.text = title_text
            _style_run(title.text_frame.paragraphs[0], 44, INK, bold=True)

            rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.15),
                                          Inches(1.5), Pt(3))
            rule.fill.solid()
            rule.fill.fore_color.rgb = ACCENT
            rule.line.fill.background()

            sub = s.get("subtitle")
            if sub and len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                ph.left, ph.top = MARGIN, Inches(4.5)
                ph.width, ph.height = W - MARGIN * 2, Inches(0.9)
                ph.text_frame.text = str(sub)
                _style_run(ph.text_frame.paragraphs[0], 18, MUTED)
            continue

        # Content slide. Layout 5 is "Title Only" — a real title placeholder
        # with no content box fighting us for the space.
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.left, title.top = MARGIN, Inches(0.62)
        title.width, title.height = W - MARGIN * 2, Inches(0.95)
        title.text_frame.text = title_text
        _style_run(title.text_frame.paragraphs[0], 30, INK, bold=True)

        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.62),
                                      Inches(0.85), Pt(3))
        rule.fill.solid()
        rule.fill.fore_color.rgb = ACCENT
        rule.line.fill.background()

        if bullets:
            size = _bullet_size(len(bullets))
            box = slide.shapes.add_textbox(MARGIN, Inches(2.05),
                                           W - MARGIN * 2, H - Inches(2.9))
            tf = box.text_frame
            tf.word_wrap = True
            for j, b in enumerate(bullets):
                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.text = b
                _style_run(para, size, BODY)
                para.space_after = Pt(size * 0.75)
                para.line_spacing = 1.25

        num = slide.shapes.add_textbox(W - MARGIN - Inches(1.2), H - Inches(0.72),
                                       Inches(1.2), Inches(0.35))
        p = num.text_frame.paragraphs[0]
        p.text = f"{i + 1} / {total}"
        p.alignment = PP_ALIGN.RIGHT
        _style_run(p, 11, MUTED)

    prs.save(path)
    return f"created {name} with {len(slides)} slide(s)"


def create_spreadsheet(files_dir, filename, rows, sheet_name=None):
    path, name = _resolve(files_dir, filename, ".xlsx")
    if not isinstance(rows, list) or not rows or not all(isinstance(r, list) for r in rows):
        raise ToolError("'rows' must be a non-empty list of row lists, e.g. "
                        '[["Item", "Cost"], ["Laptops", 3200]]')
    wb = Workbook()
    ws = wb.active
    if sheet_name:
        ws.title = str(sheet_name)
    for r in rows:
        ws.append(r)  # strings starting with '=' become real formulas

    # Presentation only — values and formulas above are untouched, so the
    # graders and read_spreadsheet() see exactly what they saw before.
    head_fill = PatternFill("solid", fgColor="1F2430")
    thin = Side(style="thin", color="D8DCE3")
    n_cols = max((len(r) for r in rows), default=0)

    for cell in ws[1]:
        cell.font = Font(name="Calibri", bold=True, size=11, color="FFFFFF")
        cell.fill = head_fill
        cell.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[1].height = 22

    last = ws.max_row
    total_row = last if str(rows[-1][0] if rows[-1] else "").strip().lower() in (
        "total", "totals", "sum") else None

    for r in range(2, last + 1):
        for c in range(1, n_cols + 1):
            cell = ws.cell(row=r, column=c)
            cell.border = Border(bottom=thin)
            cell.font = Font(name="Calibri", size=11, bold=(r == total_row))
            if isinstance(cell.value, float):
                cell.number_format = "#,##0.00"
            elif isinstance(cell.value, int) and not isinstance(cell.value, bool):
                cell.number_format = "#,##0"

    for c in range(1, n_cols + 1):
        longest = max((len(str(r[c - 1])) for r in rows
                       if len(r) >= c and r[c - 1] is not None), default=8)
        ws.column_dimensions[get_column_letter(c)].width = min(max(longest + 4, 11), 42)

    ws.freeze_panes = "A2"  # header stays put while scrolling
    wb.save(path)
    return f"created {name} with {len(rows)} row(s)"


def read_spreadsheet(files_dir, filename):
    path, name = _resolve(files_dir, filename, ".xlsx")
    if not os.path.exists(path):
        raise ToolError(f"no spreadsheet named {name} exists yet")
    wb = load_workbook(path)
    out = []
    for ws in wb.worksheets:
        rows = [[c for c in row] for row in ws.iter_rows(values_only=True)]
        out.append({"sheet": ws.title, "rows": rows})
    return out
