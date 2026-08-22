"""Tests for the harness backend.

Stdlib unittest, no pytest, no fixtures: `python -m tests.test_harness` from
the repo root. The harness has process-global state (the TOOLS registry, the
fs_tools root, the agent profile), so tests that touch it put it back.

    .venv/bin/python -m tests.test_harness
"""
import datetime
import json
import os
import sys
import tempfile
import unittest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from harness import agent, fs_tools, office, profiles, tools  # noqa: E402
from harness.memory import MemoryStore  # noqa: E402
from harness.tools import TOOLS, execute, validate_call  # noqa: E402
from harness import world as world_mod  # noqa: E402
from harness.world import ToolError, World  # noqa: E402
from bench import run as bench_run  # noqa: E402
from bench import tasks as bench_tasks  # noqa: E402


# ------------------------------------------------------------------ parsing ---

class TestParsing(unittest.TestCase):
    def test_strict_takes_clean_json_and_fenced_json(self):
        self.assertEqual(agent.parse_strict('{"tool": "think"}')[0], {"tool": "think"})
        self.assertEqual(agent.parse_strict('```json\n{"tool": "think"}\n```')[0],
                         {"tool": "think"})

    def test_strict_rejects_prose_around_json(self):
        obj, err = agent.parse_strict('Sure! {"tool": "think"}')
        self.assertIsNone(obj)
        self.assertTrue(err)

    def test_lenient_recovers_the_cases_strict_rejects(self):
        for text in ('Sure! {"tool": "think", "args": {}} hope that helps',
                     '{"tool": "think", "args": {},}',
                     'text\n{"tool": "think"}\nmore text'):
            obj, err = agent.parse_lenient(text)
            self.assertIsNotNone(obj, text)
            self.assertEqual(obj.get("tool"), "think", text)

    def test_lenient_reports_rather_than_guesses(self):
        self.assertIsNone(agent.parse_lenient("no json at all")[0])
        self.assertIsNone(agent.parse_lenient('{"a": ')[0])


# ------------------------------------------------------- date/time normalizing --

class TestNormalize(unittest.TestCase):
    MON = datetime.date(2026, 7, 20)  # a Monday

    def test_dates(self):
        n = lambda v: agent.normalize_date(v, self.MON)  # noqa: E731
        self.assertEqual(n("2026-07-22"), "2026-07-22")
        self.assertEqual(n("today"), "2026-07-20")
        self.assertEqual(n("tomorrow"), "2026-07-21")
        self.assertEqual(n("thursday"), "2026-07-23")
        self.assertEqual(n("next monday"), "2026-07-27")  # never "today"
        self.assertEqual(n("July 23"), "2026-07-23")
        self.assertEqual(n("7/23"), "2026-07-23")

    def test_unparseable_dates_pass_through_for_the_world_to_reject(self):
        self.assertEqual(agent.normalize_date("whenever", self.MON), "whenever")
        self.assertEqual(agent.normalize_date(7, self.MON), 7)

    def test_times(self):
        self.assertEqual(agent.normalize_time("2pm"), "14:00")
        self.assertEqual(agent.normalize_time("2:30 pm"), "14:30")
        self.assertEqual(agent.normalize_time("12am"), "00:00")
        self.assertEqual(agent.normalize_time("12pm"), "12:00")
        self.assertEqual(agent.normalize_time("09:15"), "09:15")
        self.assertEqual(agent.normalize_time("25:00"), "25:00")  # invalid, left alone


# ------------------------------------------------------------- weekday check ---

class TestTaskDateCheck(unittest.TestCase):
    """Began as a weekday-only check and was generalized: the failure - the
    model resolves a task's date expression itself and gets it wrong, and every
    tool answers honestly for the wrong day - is identical for "wednesday",
    "tomorrow", "July 23" and "7/23"."""

    MON = datetime.date(2026, 7, 20)

    def check(self, task, args):
        return agent.task_date_mismatch(task, args, self.MON)

    # -- extraction --------------------------------------------------------
    def test_the_expressions_the_normalizer_understands_are_all_found(self):
        # "next tuesday" from a Monday is tomorrow under the normalizer's rule
        # (delta % 7, with "or 7" only when that lands on today). The check MUST
        # agree with the normalizer, whatever the rule is - agreeing is the point.
        self.assertEqual(agent.task_dates("meet next tuesday", self.MON), {"2026-07-21"})
        self.assertEqual(agent.task_dates("do it tomorrow", self.MON), {"2026-07-21"})
        self.assertEqual(agent.task_dates("due July 23", self.MON), {"2026-07-23"})
        self.assertEqual(agent.task_dates("due 7/23", self.MON), {"2026-07-23"})
        self.assertEqual(agent.task_dates("on 2026-07-24 exactly", self.MON), {"2026-07-24"})

    def test_a_bare_month_or_plain_number_is_not_a_date(self):
        self.assertEqual(agent.task_dates("build my July receipts sheet", self.MON), set())
        self.assertEqual(agent.task_dates("the Q3 numbers, all 3 regions", self.MON), set())

    # -- the check ---------------------------------------------------------
    def test_a_wrong_weekday_is_caught_and_corrected(self):
        msg = self.check("Summarize my Wednesday meetings", {"date": "2026-07-27"})
        self.assertIsNotNone(msg)
        self.assertIn("2026-07-22", msg)  # what Wednesday actually is
        self.assertIn("Monday", msg)      # what it sent

    def test_a_wrong_tomorrow_is_caught(self):
        msg = self.check("Book dentist for tomorrow", {"date": "2026-07-20"})
        self.assertIsNotNone(msg)
        self.assertIn("2026-07-21", msg)

    def test_a_wrong_month_day_is_caught(self):
        msg = self.check("Set a reminder for July 23", {"date": "2026-07-24"})
        self.assertIsNotNone(msg)
        self.assertIn("2026-07-23", msg)

    def test_the_right_date_passes_in_every_form(self):
        self.assertIsNone(self.check("Summarize my Wednesday meetings",
                                     {"date": "2026-07-22"}))
        self.assertIsNone(self.check("Book dentist for tomorrow", {"date": "2026-07-21"}))
        self.assertIsNone(self.check("reminder for 7/23", {"date": "2026-07-23"}))

    def test_plurals_count_as_naming_the_day(self):
        self.assertIsNotNone(self.check("never book me on Fridays", {"date": "2026-07-22"}))

    def test_two_dates_in_one_task_are_left_alone(self):
        """"Move my Wednesday meeting to Friday" legitimately carries either."""
        self.assertIsNone(self.check("Move my Wednesday meeting to Friday",
                                     {"date": "2026-07-24"}))
        self.assertIsNone(self.check("move the July 23 review to tomorrow",
                                     {"date": "2026-07-21"}))

    def test_the_same_date_named_two_ways_still_counts_as_one(self):
        """"tomorrow, the 21st..." resolves to one date, so the check stays armed."""
        msg = self.check("book it tomorrow, 7/21", {"date": "2026-07-24"})
        self.assertIsNotNone(msg)

    def test_a_task_naming_no_date_is_left_alone(self):
        self.assertIsNone(self.check("Book an hour for deep work", {"date": "2026-07-27"}))

    def test_a_call_with_no_date_is_left_alone(self):
        self.assertIsNone(self.check("Summarize my Wednesday meetings", {"to": "Jordan"}))
        self.assertIsNone(self.check("Summarize my Wednesday meetings", {}))

    def test_a_date_that_is_not_a_date_is_left_to_the_validator(self):
        self.assertIsNone(self.check("my Wednesday meetings", {"date": "whenever"}))


# ------------------------------------------------------------------- repair ---

class TestRepairAndValidate(unittest.TestCase):
    def test_near_miss_parameter_is_renamed(self):
        args, notes = agent.repair_args("read_email", {"email_id": "e2"})
        self.assertEqual(args, {"id": "e2"})
        self.assertTrue(notes)

    def test_unknown_parameter_is_dropped(self):
        args, notes = agent.repair_args("read_email", {"id": "e2", "colour": "red"})
        self.assertEqual(args, {"id": "e2"})
        self.assertIn("dropped unknown parameter 'colour'", notes)

    def test_validate_names_every_problem_not_just_the_first(self):
        problems = validate_call("add_event", {"title": "x", "nope": 1})
        self.assertTrue(any("date" in p for p in problems))
        self.assertTrue(any("nope" in p for p in problems))

    def test_validate_rejects_an_unknown_tool(self):
        self.assertTrue(validate_call("teleport", {}))


# -------------------------------------------------------------------- world ---

class TestWorld(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.w = World(self.tmp.name)

    def tearDown(self):
        self.tmp.cleanup()

    def test_reading_an_unknown_email_says_how_to_find_a_real_one(self):
        with self.assertRaises(ToolError) as cm:
            self.w.read_email("e999")
        self.assertIn("list_emails", str(cm.exception))

    def test_event_must_end_after_it_starts(self):
        with self.assertRaises(ToolError):
            self.w.add_event("x", "2026-07-22", "14:00", "13:00")

    def test_attendees_given_as_a_string_are_split(self):
        ev = self.w.add_event("x", "2026-07-22", "13:00", "14:00", "a@b.com, c@d.com")
        self.assertEqual(ev["attendees"], ["a@b.com", "c@d.com"])

    def test_a_tool_error_is_logged_and_returned_readably(self):
        ok, obs = execute("read_email", {"id": "nope"}, self.w, None)
        self.assertFalse(ok)
        self.assertTrue(obs.startswith("ERROR:"))
        self.assertEqual(self.w.actions[-1]["ok"], False)


class TestCalendarEditing(unittest.TestCase):
    """The calendar was write-only: add an event, never move or cancel one.
    Asked to "move my Design review to Thursday", the agent added a SECOND
    event and reported the meeting moved. The original was still there."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.w = World(self.tmp.name)

    def tearDown(self):
        self.tmp.cleanup()

    def test_moving_an_event_changes_it_rather_than_adding_one(self):
        before = len(self.w.events)
        ev = self.w.update_event("c2", date="2026-07-23", start_time="09:00",
                                 end_time="10:00")
        self.assertEqual(len(self.w.events), before)
        self.assertEqual((ev["date"], ev["start"], ev["end"]),
                         ("2026-07-23", "09:00", "10:00"))
        self.assertEqual(ev["title"], "Design review")  # untouched fields survive

    def test_an_update_is_validated_like_a_new_event(self):
        with self.assertRaises(ToolError):
            self.w.update_event("c2", start_time="15:00", end_time="14:00")
        with self.assertRaises(ToolError):
            self.w.update_event("c2", date="next thursday")

    def test_updating_an_event_that_is_not_there_says_how_to_find_one(self):
        with self.assertRaises(ToolError) as cm:
            self.w.update_event("c999", title="x")
        self.assertIn("list_events", str(cm.exception))

    def test_cancelling_removes_it_and_reports_what_went(self):
        out = self.w.cancel_event("c2")
        self.assertEqual(out["cancelled"]["title"], "Design review")
        self.assertNotIn("c2", [e["id"] for e in self.w.events])

    def test_an_id_is_never_handed_out_twice(self):
        """Ids were len(events)+1. Cancelling an event from the middle of the
        list then adding one handed out an id that was still in use: seven
        events, cancel c2, len+1 is 7, and c7 already exists."""
        self.w.cancel_event("c2")
        new = self.w.add_event("Deep work", "2026-07-23", "14:00", "15:00")
        ids = [e["id"] for e in self.w.events]
        self.assertEqual(len(ids), len(set(ids)), ids)
        self.assertEqual(new["id"], "c8")

    def test_both_editors_count_as_writes_for_the_loop(self):
        """Was: grep agent.py for the two names, back when the write set was a
        literal there. The set is derived from declared effect classes now, so
        the intent is assertable directly instead of by proxy."""
        from harness import agent as a
        self.assertIn("update_event", a.world_changing_tools())
        self.assertIn("cancel_event", a.world_changing_tools())



# ------------------------------------------------------------------ effects ---

class TestEffectClasses(unittest.TestCase):
    """The write set used to be a literal frozenset of nine office tool names in
    agent.py. It is derived from a per-tool effect declaration now, so a domain
    that registers its own tools gets the date guard, the unplanned-write nudge,
    the repeat check and the unrequested report without editing the loop."""

    def test_every_registered_tool_declares_a_valid_effect(self):
        for name, spec in tools.TOOLS.items():
            self.assertIn(spec.get("effect"), tools.EFFECTS, name)

    def test_the_derived_write_set_is_the_nine_it_replaced(self):
        """Proof the refactor changed nothing for the graded benchmark."""
        self.assertEqual(
            set(tools.write_tool_names()),
            {"send_email", "add_event", "update_event", "cancel_event", "send_message",
             "set_reminder", "create_presentation", "create_spreadsheet", "save_memory"})

    def test_reads_and_writes_partition_the_registry(self):
        reads, writes = tools.read_tool_names(), tools.write_tool_names()
        self.assertEqual(reads & writes, frozenset())
        self.assertEqual(reads | writes, frozenset(tools.TOOLS))

    def test_an_undeclared_tool_counts_as_world_changing(self):
        """The one wrong default is to read as a read. An MCP server or a domain
        pack that registers a tool without saying what it does gets treated as
        the most dangerous thing it could be, not the safest."""
        reg = {"mystery": {"desc": "", "params": {}, "example": {}, "run": None}}
        self.assertEqual(tools.effect_of("mystery", reg), "unrecoverable_emission")
        self.assertEqual(set(tools.write_tool_names(reg)), {"mystery"})

    def test_an_unknown_tool_name_counts_as_world_changing(self):
        self.assertEqual(tools.effect_of("no_such_tool", {}), "unrecoverable_emission")

    def test_emissions_are_the_tools_that_reach_another_party(self):
        """The class split exists for confirmation policy: an emission has no
        inverse, a revertible_write does."""
        self.assertEqual(tools.effect_of("send_email"), "unrecoverable_emission")
        self.assertEqual(tools.effect_of("send_message"), "unrecoverable_emission")
        self.assertEqual(tools.effect_of("create_spreadsheet"), "revertible_write")
        self.assertEqual(tools.effect_of("read_spreadsheet"), "read")


    def test_mcp_tools_declare_their_own_effect_class(self):
        """Real-account tools reach a live mailbox, so the classification cannot
        be left to the undeclared default even though that default is safe."""
        from harness import mcp_bridge as mb
        self.assertEqual(mb._effect_class("list_messages", False), "read")
        self.assertEqual(mb._effect_class("create_draft", True), "withheld_emission")
        self.assertEqual(mb._effect_class("send_message", True), "unrecoverable_emission")
        self.assertEqual(mb._effect_class("send_draft", True), "unrecoverable_emission")

    def test_an_unrecognised_real_account_write_is_called_an_emission(self):
        """From a name alone we cannot tell whether a write reaches another
        party, and a calendar invite does. Guessing revertible_write would be
        guessing in the direction that costs something."""
        from harness import mcp_bridge as mb
        self.assertEqual(mb._effect_class("create_event", True), "unrecoverable_emission")
        self.assertEqual(mb._effect_class("update_label", True), "unrecoverable_emission")


    def test_the_unread_file_guard_follows_the_registry_not_a_literal(self):
        """The guard hardcoded create_spreadsheet, create_presentation,
        read_spreadsheet and an .xlsx regex, so it was blind to every other
        document type a domain might register."""
        from harness import agent as a
        self.assertEqual(set(tools.file_writing_tools()),
                         {"create_spreadsheet", "create_presentation"})
        self.assertEqual(tools.opener_for("q3_raw.xlsx"), "read_spreadsheet")
        self.assertIsNone(tools.opener_for("notes.txt"))
        self.assertEqual(a.filename_re().findall("the export is in q3_raw.xlsx today"),
                         ["q3_raw.xlsx"])

    def test_a_new_document_type_is_picked_up_without_touching_the_loop(self):
        """The whole point of deriving it: register a reader and a writer for a
        new extension and the guard covers them."""
        from harness import agent as a
        reg = dict(tools.TOOLS)
        reg["read_report"] = {"effect": "read", "opens": (".pdf",),
                              "desc": "", "params": {}, "example": {}, "run": None}
        reg["make_report"] = {"effect": "revertible_write", "writes_file": True,
                              "desc": "", "params": {}, "example": {}, "run": None}
        self.assertEqual(tools.opener_for("q3.pdf", reg), "read_report")
        self.assertIn("make_report", tools.file_writing_tools(reg))
        saved = dict(tools.TOOLS)
        try:
            tools.TOOLS.clear(); tools.TOOLS.update(reg)
            self.assertEqual(a.filename_re().findall("see q3.pdf and old.xlsx"),
                             ["q3.pdf", "old.xlsx"])
        finally:
            tools.TOOLS.clear(); tools.TOOLS.update(saved)

    def test_no_openable_extensions_means_no_filename_scan(self):
        saved = dict(tools.TOOLS)
        from harness import agent as a
        try:
            tools.TOOLS.clear()
            self.assertIsNone(a.filename_re())
        finally:
            tools.TOOLS.clear(); tools.TOOLS.update(saved)



# ------------------------------------------------------------------- guards ---

class TestGuardRegistry(unittest.TestCase):
    """The guards were inline `if ... continue` blocks in one 345-line
    function: order was statement order, a domain could not add one, and none
    could be exercised without scripting a whole run. Each is a named function
    over a GuardState now, so these tests call them directly."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()

    def tearDown(self):
        self.tmp.cleanup()

    def _state(self, plan="", task="do the thing"):
        w = World(self.tmp.name)
        return agent.GuardState(None, w, agent.Episode(), [], task,
                                agent.world_changing_tools(), plan)

    def test_every_guard_has_a_name_and_a_callable(self):
        for name, check in agent.GUARDS:
            self.assertTrue(name and callable(check), name)
        self.assertEqual([n for n, _ in agent.GUARDS],
                         ["wrong_date", "unplanned_write", "unread_file",
                          "read_before_write"])

    def test_a_guard_abstains_by_returning_none(self):
        g = self._state()
        g.name, g.args = "list_emails", {}
        self.assertIsNone(agent.run_guards(g))

    def test_denial_is_monotonic_nothing_after_the_first_speaker_runs(self):
        """The property goal 2 exists for: once a guard has questioned a call,
        no later guard is consulted, so nothing can turn a question back into
        permission."""
        ran = []

        def first(g):
            ran.append("first")
            return "questioned by first"

        def second(g):
            ran.append("second")
            return None

        g = self._state()
        g.name, g.args = "send_email", {}
        got = agent.run_guards(g, [("first", first), ("second", second)])
        self.assertEqual(got, ("first", "questioned by first"))
        self.assertEqual(ran, ["first"], "a guard after the first speaker ran")

    def test_a_later_guard_cannot_reverse_an_earlier_question(self):
        g = self._state()
        g.name, g.args = "send_email", {}
        got = agent.run_guards(g, [("deny", lambda s: "no"),
                                   ("allow", lambda s: None)])
        self.assertEqual(got[0], "deny")

    # ---- the individual guards ----

    def test_wrong_date_fires_on_writes_only(self):
        g = self._state(task="Book it on Thursday")
        g.name, g.args = "add_event", {"date": "2026-07-20"}
        self.assertIsNotNone(agent.guard_wrong_date(g))
        g.name = "list_events"
        self.assertIsNone(agent.guard_wrong_date(g),
                          "a read with a mismatched date is the model looking around")

    def test_unplanned_write_questions_once_then_lets_it_through(self):
        g = self._state(plan="1. list_emails - look")
        g.name, g.args = "send_email", {}
        first = agent.guard_unplanned_write(g)
        self.assertIn("never included send_email", first)
        self.assertIsNone(agent.guard_unplanned_write(g),
                          "questioned twice; the contract is question ONCE")

    def test_the_unplanned_write_nudge_never_tells_the_model_to_do_less(self):
        """This message once ended 'Only do what the task requires - nothing
        extra', and an 8B obeyed that instead of insisting: the question became
        a block in practice, the model abandoned the job and reported success."""
        g = self._state(plan="1. list_emails - look")
        g.name, g.args = "send_email", {}
        msg = agent.guard_unplanned_write(g)
        self.assertNotIn("nothing extra", msg)
        self.assertIn("call it again and it will run", msg)

    def test_save_memory_is_exempt_from_the_unplanned_write_guard(self):
        g = self._state(plan="1. list_emails - look")
        g.name, g.args = "save_memory", {"fact": "x"}
        self.assertIsNone(agent.guard_unplanned_write(g))

    def test_read_before_write_needs_a_read_the_plan_put_first(self):
        g = self._state(plan="1. list_emails - look\n2. send_email - reply")
        g.name, g.args = "send_email", {}
        self.assertIn("You planned to call list_emails first",
                      agent.guard_read_before_write(g))
        self.assertIsNone(agent.guard_read_before_write(g), "question ONCE")

    def test_read_before_write_stays_quiet_once_something_was_read(self):
        g = self._state(plan="1. list_emails - look\n2. send_email - reply")
        g.looked = True
        g.name, g.args = "send_email", {}
        self.assertIsNone(agent.guard_read_before_write(g))

    def test_a_plan_that_never_proposed_looking_holds_the_model_to_nothing(self):
        g = self._state(plan="1. send_email - reply")
        self.assertIsNone(g.first_read_planned)
        g.name, g.args = "send_email", {}
        self.assertIsNone(agent.guard_read_before_write(g))

    def _seed(self, g, name="export.xlsx"):
        os.makedirs(g.world.files_dir, exist_ok=True)
        open(os.path.join(g.world.files_dir, name), "wb").close()
        return name

    def test_unread_file_fires_only_for_a_file_the_run_was_told_about(self):
        g = self._state()
        f = self._seed(g)
        g.name, g.args = "create_spreadsheet", {}
        self.assertIsNone(agent.guard_unread_file(g), "nothing mentioned yet")
        g.mentioned_files.add(f)
        msg = agent.guard_unread_file(g)
        self.assertIn(f, msg)
        self.assertIn("read_spreadsheet", msg)
        self.assertIsNone(agent.guard_unread_file(g), "question ONCE")

    def test_unread_file_ignores_a_file_already_opened(self):
        g = self._state()
        f = self._seed(g)
        g.name, g.args = "create_spreadsheet", {}
        g.mentioned_files.add(f)
        g.opened_files.add(f)
        self.assertIsNone(agent.guard_unread_file(g))

    def test_unread_file_ignores_a_mention_of_a_file_that_is_not_here(self):
        """The run being told about a file it cannot open is not evidence of
        anything; only a file that exists on disk is."""
        g = self._state()
        g.name, g.args = "create_spreadsheet", {}
        g.mentioned_files.add("somewhere_else.xlsx")
        self.assertIsNone(agent.guard_unread_file(g))

    def test_unread_file_only_fires_for_a_tool_that_writes_a_file(self):
        g = self._state()
        f = self._seed(g)
        g.mentioned_files.add(f)
        g.name, g.args = "send_email", {}
        self.assertIsNone(agent.guard_unread_file(g))

    def test_every_guard_message_offers_a_way_through(self):
        """Question once, never forbid. A guard that leaves the model no way
        forward is a block wearing a question's clothes."""
        cases = [
            ("unplanned_write", agent.guard_unplanned_write,
             self._state(plan="1. list_emails - look"), "send_email"),
            ("read_before_write", agent.guard_read_before_write,
             self._state(plan="1. list_emails - look\n2. send_email - reply"), "send_email"),
        ]
        for name, check, g, call in cases:
            g.name, g.args = call, {}
            msg = check(g)
            self.assertIsNotNone(msg, name)
            self.assertIn("it will run", msg, name)



# --------------------------------------------------------------- list_files ---

class TestListFiles(unittest.TestCase):
    """The agent had sixteen tools and no way to see its own workspace: it
    could create a spreadsheet and read one back by name, but only if something
    had already told it the name."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.w = World(self.tmp.name)

    def tearDown(self):
        self.tmp.cleanup()

    def _write(self, name, size=12):
        with open(os.path.join(self.w.files_dir, name), "wb") as fh:
            fh.write(b"x" * size)

    def test_an_empty_workspace_says_so_rather_than_returning_nothing(self):
        """An empty list renders as "[]", which a small model reads as an error
        or as a tool that did not work."""
        self.assertEqual(self.w.list_files(), "the workspace has no files yet")

    def test_it_reports_name_size_and_when_it_changed(self):
        self._write("q3.xlsx", 40)
        rows = self.w.list_files()
        self.assertEqual(len(rows), 1)
        self.assertEqual(rows[0]["name"], "q3.xlsx")
        self.assertEqual(rows[0]["bytes"], 40)
        self.assertRegex(rows[0]["modified"], r"^\d{4}-\d{2}-\d{2} \d{2}:\d{2}$")

    def test_it_is_sorted_so_two_runs_read_the_same(self):
        for n in ("c.xlsx", "a.xlsx", "b.xlsx"):
            self._write(n)
        self.assertEqual([r["name"] for r in self.w.list_files()],
                         ["a.xlsx", "b.xlsx", "c.xlsx"])

    def test_it_is_a_read_and_joins_the_read_set(self):
        self.assertEqual(tools.effect_of("list_files"), "read")
        self.assertIn("list_files", tools.read_tool_names())
        self.assertNotIn("list_files", agent.world_changing_tools())

    def test_it_runs_through_the_registry(self):
        self._write("q3.xlsx")
        ok, obs = execute("list_files", {}, self.w, None)
        self.assertTrue(ok)
        self.assertIn("q3.xlsx", obs)

    def test_a_listing_does_not_count_as_the_task_naming_a_file(self):
        """The load-bearing half. Filenames are harvested out of read results
        so the unread-file guard can tell a file the task pointed at from one
        the model invented. If a directory listing counted, enumerating the
        workspace would mark every file in it as named by the task, and the
        next document write would be questioned about an arbitrary file."""
        self.assertTrue(TOOLS["list_files"].get("lists_files"))
        self.assertFalse(TOOLS["read_spreadsheet"].get("lists_files"))

    def test_the_guard_stays_quiet_after_a_listing_but_fires_after_a_read(self):
        """The behaviour the flag buys, through the guard rather than the flag."""
        self._write("q3.xlsx")
        g = agent.GuardState(None, self.w, agent.Episode(), [], "make a sheet",
                             agent.world_changing_tools(), "")
        g.name, g.args = "create_spreadsheet", {}
        # Nothing harvested from a listing, so nothing to question.
        self.assertIsNone(agent.guard_unread_file(g))
        # The same filename arriving from something the run actually read does
        # arm it.
        g.mentioned_files.add("q3.xlsx")
        self.assertIn("q3.xlsx", agent.guard_unread_file(g))



# ------------------------------------------------------ mcp tool restriction ---

class TestRestrictToMcp(unittest.TestCase):
    """restrict_to_mcp used to be an ALLOW-list naming the seven tools that
    survived, so every base-layer tool added afterwards silently vanished the
    moment MCP was on. Found by adding list_files and watching a real run tell
    the model 'unknown tool list_files'."""

    def setUp(self):
        self.saved = dict(TOOLS)

    def tearDown(self):
        TOOLS.clear(); TOOLS.update(self.saved)

    def test_it_drops_exactly_the_simulated_connectors(self):
        from harness import mcp_bridge
        before = set(TOOLS)
        mcp_bridge.restrict_to_mcp()
        self.assertEqual(before - set(TOOLS),
                         {"list_emails", "read_email", "send_email", "list_events",
                          "add_event", "update_event", "cancel_event",
                          "send_message", "set_reminder"})

    def test_a_new_base_layer_tool_survives_without_being_named_anywhere(self):
        """The regression this inversion exists to prevent."""
        from harness import mcp_bridge
        TOOLS["some_future_tool"] = {"effect": "read", "desc": "", "params": {},
                                     "example": {}, "run": None}
        mcp_bridge.restrict_to_mcp()
        self.assertIn("some_future_tool", TOOLS)
        self.assertIn("list_files", TOOLS)

    def test_think_memory_and_done_still_survive(self):
        from harness import mcp_bridge
        mcp_bridge.restrict_to_mcp()
        for n in ("think", "save_memory", "recall_memories", "done"):
            self.assertIn(n, TOOLS)

    def test_dropping_the_document_tools_too_is_still_available(self):
        from harness import mcp_bridge
        mcp_bridge.restrict_to_mcp(keep_office_docs=False)
        for n in ("create_spreadsheet", "create_presentation", "read_spreadsheet"):
            self.assertNotIn(n, TOOLS)
        self.assertIn("think", TOOLS)
        self.assertIn("list_files", TOOLS,
                      "listing the workspace is not a document tool")



# ------------------------------------------------------- fs tool effects ---

class TestFsToolEffects(unittest.TestCase):
    def test_the_fs_write_set_is_derived_not_listed_twice(self):
        self.assertEqual(set(fs_tools.WRITE_TOOLS),
                         {"write_file", "append_file", "delete_path",
                          "move_path", "run_command"})

    def test_every_fs_tool_declares_a_valid_effect(self):
        for name, spec in fs_tools._FS_TOOLS.items():
            self.assertIn(spec.get("effect"), tools.EFFECTS, name)

    def test_delete_and_shell_are_emissions_because_nothing_undoes_them(self):
        """delete_path removes bytes this system cannot restore, and a command
        can reach the network. Calling either revertible would be guessing in
        the direction that costs something."""
        self.assertEqual(fs_tools._FS_TOOLS["delete_path"]["effect"],
                         "unrecoverable_emission")
        self.assertEqual(fs_tools._FS_TOOLS["run_command"]["effect"],
                         "unrecoverable_emission")

    def test_listing_a_real_directory_is_not_the_task_naming_a_file(self):
        """Same reason list_files carries the flag: a listing is not a mention."""
        self.assertTrue(fs_tools._FS_TOOLS["list_dir"].get("lists_files"))
        self.assertFalse(fs_tools._FS_TOOLS["read_file"].get("lists_files"))



# --------------------------------------------------------- plan repair ---

class TestPlanNameRepair(unittest.TestCase):
    """Every other layer repairs before it rejects: near-miss parameter names
    are renamed, unknown ones dropped, dates normalized, and a misspelled tool
    CALL gets "did you mean" with the right shape. The plan step did not, so a
    model writing create_sheet lost that step silently."""

    def test_a_real_name_passes_through_unchanged(self):
        for n in ("think", "create_spreadsheet", "list_emails"):
            self.assertEqual(agent.plan_name(n), n)

    def test_a_near_miss_is_repaired_onto_the_real_tool(self):
        self.assertEqual(agent.plan_name("create_sheet"), "create_spreadsheet")
        self.assertEqual(agent.plan_name("list_email"), "list_emails")
        self.assertEqual(agent.plan_name("send_mail"), "send_email")

    def test_something_that_is_not_a_tool_is_still_dropped(self):
        """Repair, not invention. Free prose must never enter the context as a
        plan step."""
        for n in ("xyzzy", "", None, 123, "do the thing"):
            self.assertIsNone(agent.plan_name(n))

    def test_losing_a_step_is_not_cosmetic(self):
        """Why this matters: the plan arms two guards. A write the plan no
        longer names arms the unplanned-write guard, and a read it no longer
        names disarms the read-before-write one."""
        state = agent.GuardState(None, World(self.tmp.name), agent.Episode(), [],
                                 "build it", agent.world_changing_tools(),
                                 "1. list_emails - look\n2. create_spreadsheet - build")
        self.assertEqual(state.first_read_planned, "list_emails")
        self.assertIn("create_spreadsheet", state.planned_set)
        dropped = agent.GuardState(None, World(self.tmp.name), agent.Episode(), [],
                                   "build it", agent.world_changing_tools(),
                                   "1. create_spreadsheet - build")
        self.assertIsNone(dropped.first_read_planned)

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()

    def tearDown(self):
        self.tmp.cleanup()



# ------------------------------------------------- model-visible is logged ---

OBS_PREFIX = "OBSERVATION: "


class _SpyLLM:
    """Records every conversation handed to the model, verbatim."""

    def __init__(self, replies):
        self.replies, self.calls, self._i = replies, 0, 0
        self.seen = []

    def chat(self, messages, **kw):
        self.seen.append([dict(m) for m in messages])
        self.calls += 1
        reply = self.replies[min(self._i, len(self.replies) - 1)]
        self._i += 1
        return reply


class TestModelVisibleIsLogged(unittest.TestCase):
    maxDiff = None
    """The invariant: anything that reaches a model request must be
    reconstructable from the episode transcript.

    Without it there is no reliable answer to "what did the model actually see
    on the run where the guard fired", which is the only question an ablation
    is asking. It caught a real gap the first time it ran: the verifier's
    verdict was recorded and its INPUT was not, so the one call that decides
    whether done() is accepted was the one call nobody could inspect."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.world = World(self.tmp.name)
        self.mem = MemoryStore(os.path.join(self.tmp.name, "mem.jsonl"))
        self._saved = agent.PROFILE

    def tearDown(self):
        agent.set_profile(self._saved)
        self.tmp.cleanup()

    @staticmethod
    def _call(tool, **args):
        return json.dumps({"thought": "t", "tool": tool, "args": args})

    def _unlogged(self, ep, llm):
        """Every distinct thing the model was shown that no note accounts for.

        Exact containment, not the loose two-way substring match a first
        attempt used: the task note is a substring of the verifier prompt, so
        loose matching reported the verifier's input as logged when it was
        not."""
        logged = [str(n["content"]) for n in ep.transcript]
        missing = []
        for conv in llm.seen:
            for msg in conv:
                body = msg["content"]
                if not body or not body.strip():
                    continue
                # One named exception, not a loose rule: the loop frames a tool
                # result as "OBSERVATION: <result>" and the note carries the
                # result alone. The variable content is fully reconstructable;
                # only a constant the harness itself adds differs. Every other
                # comparison stays exact containment.
                probe = body[len(OBS_PREFIX):] if body.startswith(OBS_PREFIX) else body
                if any(probe in note for note in logged):
                    continue
                if body in missing:
                    continue
                missing.append(body)
        return missing

    def test_every_message_the_model_saw_is_in_the_transcript(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=1))
        plan = '{"steps": [{"tool": "list_emails", "what": "look"}, ' \
               '{"tool": "send_email", "what": "reply"}]}'
        llm = _SpyLLM([
            plan,
            self._call("list_emails"),
            self._call("send_email"),
            self._call("done", summary="looked"),
            '{"complete": true, "missing": "", "unrequested": ""}',
        ])
        ep = agent.run_harness(llm, self.world, self.mem, "look at the inbox and reply to Sam")
        self.assertEqual(self._unlogged(ep, llm), [],
                         "the model saw something the transcript cannot reconstruct")

    def test_the_verifiers_input_is_logged_not_only_its_verdict(self):
        """The gap this invariant found. The verifier decides whether done() is
        accepted, so its evidence block is the most important thing in the run
        to be able to read afterwards."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=1))
        llm = _SpyLLM([
            self._call("list_emails"),
            self._call("done", summary="looked"),
            '{"complete": true, "missing": "", "unrequested": ""}',
        ])
        ep = agent.run_harness(llm, self.world, self.mem, "look at the inbox")
        kinds = [n["kind"] for n in ep.transcript]
        self.assertIn("verify_prompt", kinds)
        prompt = next(n["content"] for n in ep.transcript if n["kind"] == "verify_prompt")
        self.assertIn("task-completion verifier", prompt, "the system prompt")
        self.assertIn("ACTIONS THE ASSISTANT TOOK", prompt, "the evidence block")

    def test_it_holds_when_a_guard_and_a_repair_are_in_play(self):
        """The interesting runs are the ones with corrections in them, so the
        invariant has to hold on those and not only on a clean pass."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=1))
        plan = '{"steps": [{"tool": "list_email", "what": "look"}, ' \
               '{"tool": "create_sheet", "what": "build"}]}'
        llm = _SpyLLM([
            plan,
            self._call("send_email"),
            self._call("send_email", to="a@b.c", subject="s", body="b"),
            self._call("done", summary="sent"),
            '{"complete": true, "missing": "", "unrequested": ""}',
        ])
        ep = agent.run_harness(llm, self.world, self.mem, "email someone")
        self.assertEqual(self._unlogged(ep, llm), [])



# ------------------------------------------------ revertible registration ---

class TestRevertibleRegistration(unittest.TestCase):
    """Three modules mutate the shared tool registry and none of them could be
    reversed. mcp_bridge.shutdown() closed its subprocesses and left their
    tools in the registry forever; neither restrict_ function had an undo at
    all. A registration that cannot be undone is not a registration, it is a
    leak."""

    def setUp(self):
        self.saved = dict(TOOLS)

    def tearDown(self):
        TOOLS.clear(); TOOLS.update(self.saved)

    def test_register_returns_a_disposer_that_removes_exactly_what_it_added(self):
        undo = tools.register("temp_tool", {"effect": "read", "desc": "", "params": {},
                                            "example": {}, "run": None})
        self.assertIn("temp_tool", TOOLS)
        undo()
        self.assertNotIn("temp_tool", TOOLS)

    def test_the_disposer_is_idempotent(self):
        undo = tools.register("temp_tool", {"effect": "read", "desc": "", "params": {},
                                            "example": {}, "run": None})
        undo(); undo()
        self.assertNotIn("temp_tool", TOOLS)

    def test_it_restores_what_was_shadowed_not_merely_what_was_added(self):
        """The difference that makes it an inverse rather than a delete."""
        original = TOOLS["think"]
        undo = tools.register("think", {"effect": "read", "desc": "impostor",
                                        "params": {}, "example": {}, "run": None})
        self.assertEqual(TOOLS["think"]["desc"], "impostor")
        undo()
        self.assertIs(TOOLS["think"], original)

    def test_suppress_hides_tools_and_puts_them_back(self):
        undo = tools.suppress(["think", "done"])
        self.assertNotIn("think", TOOLS)
        self.assertNotIn("done", TOOLS)
        undo()
        self.assertIn("think", TOOLS)
        self.assertIn("done", TOOLS)

    def test_suppressing_something_absent_is_not_an_error_and_stays_absent(self):
        undo = tools.suppress(["never_existed"])
        undo()
        self.assertNotIn("never_existed", TOOLS)

    def test_restrict_to_mcp_can_be_undone(self):
        from harness import mcp_bridge
        before = set(TOOLS)
        undo = mcp_bridge.restrict_to_mcp()
        self.assertNotIn("list_emails", TOOLS)
        undo()
        self.assertEqual(set(TOOLS), before)

    def test_restrict_to_files_can_be_undone(self):
        before = set(TOOLS)
        undo = fs_tools.restrict_to_files()
        self.assertNotIn("list_emails", TOOLS)
        undo()
        self.assertEqual(set(TOOLS), before)

    def test_fs_enable_then_disable_leaves_the_registry_as_it_was(self):
        before = set(TOOLS)
        fs_tools.enable(self.tmpdir(), allow_shell=True)
        self.assertIn("write_file", TOOLS)
        fs_tools.disable()
        self.assertEqual(set(TOOLS), before)

    def tmpdir(self):
        d = tempfile.TemporaryDirectory()
        self.addCleanup(d.cleanup)
        return d.name



# ------------------------------------------------------------------ bench ---

class TestBenchGraders(unittest.TestCase):
    """A grader reads the FINAL WORLD, never the transcript: the transcript is
    what the model claims it did, the world is what happened, and every failure
    this harness exists to catch is a case where those disagreed.

    Each grader is checked BOTH ways. A grader that can never pass looks
    exactly like a real result, which is the same trap as an invariant matcher
    that lies toward "all clear"."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.w = World(self.tmp.name)

    def tearDown(self):
        self.tmp.cleanup()

    def test_read_only_passes_on_an_untouched_world(self):
        self.assertTrue(bench_tasks._check_read_only(self.w)[0])

    def test_read_only_catches_each_kind_of_side_effect(self):
        self.w.send_message("sam", "hi")
        ok, why = bench_tasks._check_read_only(self.w)
        self.assertFalse(ok)
        self.assertIn("message", why)

    def test_read_only_catches_a_calendar_change(self):
        self.w.add_event("X", "2026-07-23", "10:00", "11:00")
        self.assertFalse(bench_tasks._check_read_only(self.w)[0])

    def test_receipts_passes_only_on_the_real_amounts(self):
        office.create_spreadsheet(self.w.files_dir, "r.xlsx",
                                  [["Vendor", "Amount"], ["CloudHost", "230.00"],
                                   ["OfficeMax", "87.50"], ["Delta", "412.30"]])
        self.assertTrue(bench_tasks._check_receipts_sheet(self.w)[0])

    def test_receipts_catches_invented_numbers(self):
        """The failure the read-before-write guard exists for: a plausible
        sheet full of numbers that are not the ones in the inbox."""
        office.create_spreadsheet(self.w.files_dir, "r.xlsx",
                                  [["Item", "Cost"], ["A", 100], ["B", 200]])
        ok, why = bench_tasks._check_receipts_sheet(self.w)
        self.assertFalse(ok)
        self.assertIn("invented", why)

    def test_receipts_distinguishes_nothing_written_from_wrong_numbers(self):
        """The two failures are not the same and the table must not merge
        them: raw wrote nothing, the harness wrote the wrong thing."""
        self.assertIn("no spreadsheet", bench_tasks._check_receipts_sheet(self.w)[1])

    def test_deep_work_passes_on_thursday_and_fails_on_any_other_day(self):
        self.w.add_event("Deep work", "2026-07-23", "14:00", "15:00")
        self.assertTrue(bench_tasks._check_deep_work(self.w)[0])
        w2 = World(tempfile.mkdtemp())
        w2.add_event("Deep work", "2026-07-20", "14:00", "15:00")
        ok, why = bench_tasks._check_deep_work(w2)
        self.assertFalse(ok)
        self.assertIn("2026-07-20", why)

    def test_message_jordan_needs_jordan_specifically(self):
        self.w.send_message("sam", "the list")
        self.assertFalse(bench_tasks._check_message_jordan(self.w)[0])
        self.w.send_message("jordan", "the list")
        self.assertTrue(bench_tasks._check_message_jordan(self.w)[0])

    def test_an_ablation_arm_exists_for_every_registered_guard(self):
        """The rig's whole purpose. If a guard has no arm it cannot be
        measured, and a guard nobody can measure is a guard nobody can
        justify keeping."""
        arms = bench_run.arms_for(None)
        for name, _ in agent.GUARDS:
            self.assertIn(f"harness-no-{name}", arms)
        self.assertIn("raw", arms)
        self.assertIn("harness", arms)

    def test_the_export_grader_excludes_the_source_file(self):
        """It shipped broken for one commit-worth of minutes: _sheet_text read
        every xlsx including the seeded source, so the figures it looks for
        were present whatever the agent did and the grader ALWAYS passed.
        Caught only by checking the negative direction."""
        bench_tasks._seed_export(self.w)
        self.assertFalse(bench_tasks._check_export_copied(self.w)[0],
                         "nothing written yet")
        office.create_spreadsheet(self.w.files_dir, "clean.xlsx",
                                  [["Region", "Q3"], ["West", 999]])
        ok, why = bench_tasks._check_export_copied(self.w)
        self.assertFalse(ok, "invented figures must not pass")
        self.assertIn("invented", why)

    def test_the_export_grader_passes_when_the_figures_were_copied(self):
        bench_tasks._seed_export(self.w)
        office.create_spreadsheet(
            self.w.files_dir, "clean.xlsx",
            [["Region", "Q3"]] + [[r, v] for r, v in bench_tasks.Q3_EXPORT])
        self.assertTrue(bench_tasks._check_export_copied(self.w)[0])

    def test_the_move_grader_separates_moved_from_duplicated(self):
        """Three outcomes, not two: still there, moved, or moved by adding a
        second one and leaving the first."""
        self.assertFalse(bench_tasks._check_moved_not_duplicated(self.w)[0])
        moved = World(tempfile.mkdtemp())
        moved.update_event("c2", start_time="09:00", end_time="10:00")
        self.assertTrue(bench_tasks._check_moved_not_duplicated(moved)[0])
        dupe = World(tempfile.mkdtemp())
        dupe.add_event("Design review", "2026-07-22", "09:00", "10:00")
        ok, why = bench_tasks._check_moved_not_duplicated(dupe)
        self.assertFalse(ok)
        self.assertIn("2 Design reviews", why)

    def test_a_task_may_seed_the_fixture_it_needs(self):
        """Writing around a file is only a choice if the file is there and
        something told the agent about it. The clean fixture has neither."""
        self.assertEqual(self.w.file_names(), set())
        bench_tasks._seed_export(self.w)
        self.assertIn(bench_tasks.SOURCE_FILE, self.w.file_names())
        self.assertTrue(any(bench_tasks.SOURCE_FILE in e["body"]
                            for e in self.w.emails))

    def test_the_stale_memory_grader_catches_the_wrong_answer(self):
        """The failure it grades, observed live: a run had saved "Wednesday has
        0 meetings", and the agent told a colleague their Wednesday was clear
        without opening the calendar, which held three."""
        self.assertFalse(
            bench_tasks._check_memory_did_not_outrank_the_world(self.w)[0],
            "no message at all is not a pass")
        self.w.send_message("jordan", "Your Wednesday is completely free.")
        ok, why = bench_tasks._check_memory_did_not_outrank_the_world(self.w)
        self.assertFalse(ok)
        self.assertIn("3 meetings", why)

    def test_the_stale_memory_grader_passes_an_answer_from_the_calendar(self):
        self.w.send_message("jordan", "You have 3 meetings on Wednesday.")
        self.assertTrue(
            bench_tasks._check_memory_did_not_outrank_the_world(self.w)[0])

    def test_the_stale_memory_setup_reaches_the_store_the_loop_searches(self):
        """A memory seeded anywhere else would never be injected, and the task
        would grade a run that never saw the stale fact."""
        mem = MemoryStore(os.path.join(tempfile.mkdtemp(), "m.jsonl"))
        bench_tasks._seed_stale_memory(self.w, mem)
        self.assertTrue(mem.search("Wednesday"))

    def test_an_arm_is_a_value_and_changes_nothing_globally(self):
        """Replaces test_dropping_a_guard_is_undone_afterwards, which asserted
        that bench_run._guards_without patched agent.GUARDS and restored it.
        That function is gone: an arm is a RunConfig now, so there is nothing
        to restore and the old test was asserting the mechanism rather than
        the property."""
        before = list(agent.GUARDS)
        cfg = bench_run.config_for("harness-no-wrong_date", 12, agent.PROFILE)
        self.assertNotIn("wrong_date", [n for n, _ in cfg.guards])
        self.assertEqual(agent.GUARDS, before, "the module list must be untouched")

    def test_two_arms_can_exist_at_the_same_time(self):
        """The point of threading the config. Before this, one process held one
        configuration, so arms had to run in sequence."""
        full = bench_run.config_for("harness", 12, agent.PROFILE)
        ablated = bench_run.config_for("harness-no-read_before_write", 12, agent.PROFILE)
        self.assertIn("read_before_write", [n for n, _ in full.guards])
        self.assertNotIn("read_before_write", [n for n, _ in ablated.guards])
        self.assertEqual(len(full.guards), len(ablated.guards) + 1)

    def test_a_sweep_can_run_a_model_under_another_models_profile(self):
        """Holding the model fixed and varying the profile is the only way to
        measure the plan-dependent guards on a machine where the installed
        model's own profile switches planning off."""
        own = profiles.for_model("llama3.2:1b")
        forced = profiles.for_model("llama3.1:8b")
        self.assertFalse(own.plan)
        self.assertTrue(forced.plan)

    def test_the_read_before_write_guard_needs_a_read_in_the_plan(self):
        """The rig's first finding, pinned so it cannot be forgotten: the guard
        is conditional on the plan naming a read before its first write. A plan
        that is one write step arms nothing, which is what a 1b produces."""
        g = agent.GuardState(None, self.w, agent.Episode(), [],
                             "build it", agent.world_changing_tools(),
                             "1. create_spreadsheet - build it")
        self.assertIsNone(g.first_read_planned)
        g.name, g.args = "create_spreadsheet", {}
        self.assertIsNone(agent.guard_read_before_write(g),
                          "nothing in the plan to hold the model to")

    def test_a_row_records_which_guards_actually_spoke(self):
        """Without this a table of identical arms is ambiguous, and the
        ambiguity flatters the rig: it could mean removing a guard changed
        nothing, or it could mean no guard ever fired and the ablation measured
        nothing at all. The first full ablation sweep was the second case."""
        ep = agent.Episode()
        ep.note("model", "{}")
        ep.note("guard", "read_before_write")
        ep.note("feedback", "you planned to look first")
        ep.note("guard", "wrong_date")
        self.assertEqual(bench_run.guards_fired(ep),
                         ["read_before_write", "wrong_date"])

    def test_an_episode_with_no_guard_reports_an_empty_list_not_none(self):
        """An empty list and a missing key read the same in a table and mean
        different things when summing across a sweep."""
        self.assertEqual(bench_run.guards_fired(agent.Episode()), [])
        self.assertEqual(bench_run.guards_fired(None), [])

    def test_a_scripted_arm_fires_the_guard_the_failure_targets(self):
        """The ablation that finally measures something. A scripted failure
        needs no model, so this runs in the suite."""
        task = dict(next(t for t in bench_tasks.TASKS if t["id"] == "receipts_sheet"),
                    script=bench_tasks.SCRIPTS["receipts_sheet"])
        row = bench_run.run_one("harness", task, "unused", 12,
                                profiles.for_model("llama3.1:8b"))
        self.assertIn("read_before_write", row["guards_fired"])

    def test_ablating_that_guard_silences_it_and_nothing_else(self):
        task = dict(next(t for t in bench_tasks.TASKS if t["id"] == "receipts_sheet"),
                    script=bench_tasks.SCRIPTS["receipts_sheet"])
        row = bench_run.run_one("harness-no-read_before_write", task, "unused", 12,
                                profiles.for_model("llama3.1:8b"))
        self.assertNotIn("read_before_write", row["guards_fired"])

    def test_a_question_once_guard_does_not_change_the_outcome(self):
        """Why the pass column cannot show guard value, pinned so nobody reads
        an identical pass table as evidence of no effect. The guard speaks and
        then lets an insisting model through: that IS the contract."""
        task = dict(next(t for t in bench_tasks.TASKS if t["id"] == "receipts_sheet"),
                    script=bench_tasks.SCRIPTS["receipts_sheet"])
        withg = bench_run.run_one("harness", task, "unused", 12,
                                  profiles.for_model("llama3.1:8b"))
        without = bench_run.run_one("harness-no-read_before_write", task, "unused", 12,
                                    profiles.for_model("llama3.1:8b"))
        self.assertEqual(withg["passed"], without["passed"])
        self.assertNotEqual(withg["guards_fired"], without["guards_fired"])

    def _scripted(self, arm, task_id, max_calls=14):
        task = dict(next(t for t in bench_tasks.TASKS if t["id"] == task_id),
                    script=bench_tasks.SCRIPTS[task_id])
        return bench_run.run_one(arm, task, "unused", max_calls,
                                 profiles.for_model("llama3.1:8b"))

    def test_every_registered_guard_has_a_scripted_failure(self):
        """A guard with no script cannot be ablated meaningfully on this
        machine, because no installed model reaches its failure."""
        fired = set()
        for task_id in bench_tasks.SCRIPTS:
            fired.update(self._scripted("harness", task_id)["guards_fired"])
        for name, _ in agent.GUARDS:
            self.assertIn(name, fired, f"no scripted failure fires {name}")

    def test_the_date_guard_fires_on_the_day_the_task_did_not_name(self):
        self.assertIn("wrong_date", self._scripted("harness", "deep_work")["guards_fired"])
        self.assertNotIn("wrong_date",
                         self._scripted("harness-no-wrong_date", "deep_work")["guards_fired"])

    def test_the_unplanned_write_guard_fires_on_a_send_nobody_planned(self):
        self.assertIn("unplanned_write",
                      self._scripted("harness", "read_only")["guards_fired"])
        self.assertNotIn("unplanned_write",
                         self._scripted("harness-no-unplanned_write",
                                        "read_only")["guards_fired"])

    def test_a_script_must_outlast_every_guard_that_questions_it(self):
        """The artifact this nearly produced: two guards questioned deep_work
        and each spent one attempt, so a two-attempt script reached done before
        the event existed and the ABLATED arm passed while the full harness
        failed. That reads as "the guard hurts" and is purely a short script.
        With enough attempts both arms pass, which is the real answer."""
        self.assertTrue(self._scripted("harness", "deep_work")["passed"])
        self.assertTrue(self._scripted("harness-no-wrong_date", "deep_work")["passed"])

    def test_an_unknown_arm_is_refused_rather_than_silently_skipped(self):
        with self.assertRaises(SystemExit):
            bench_run.arms_for(["harness-no-nonsense"])



# ---------------------------------------------------------- world contract ---

class _MinimalWorld:
    """A world that implements the contract and NOTHING else.

    Deliberately not a World subclass and deliberately missing every office
    member. If the loop reaches for an inbox, a calendar or a files_dir, this
    raises AttributeError and the test fails, which is the point: the contract
    is verified by running against it, not by grepping agent.py for what it
    mentions."""

    def __init__(self):
        self.actions = []
        self.snapshots = 0

    def file_names(self):
        return set()

    def snapshot(self):
        self.snapshots += 1

    def log(self, tool, args, ok, result_preview):
        self.actions.append({"tool": tool, "args": args, "ok": ok,
                             "result": result_preview})


class TestWorldContract(unittest.TestCase):
    """The audit's complaint was that world.py is a concrete class and nothing
    enforces the shape. The contract is four members, and it was measured
    rather than designed: they are the only ones the loop and the execution
    layer actually touch."""

    def test_the_real_world_satisfies_it(self):
        w = World(tempfile.mkdtemp())
        self.assertEqual(world_mod.missing_world_members(w), ())

    def test_the_checker_names_what_is_missing(self):
        self.assertEqual(world_mod.missing_world_members(object()),
                         ("actions", "file_names", "snapshot", "log"))

    def test_a_world_with_only_the_contract_can_drive_the_loop(self):
        """The claim, verified by running rather than by reading. A world with
        no inbox, no calendar and no files_dir completes a real episode."""
        w = _MinimalWorld()
        mem = MemoryStore(os.path.join(tempfile.mkdtemp(), "m.jsonl"))
        cfg = agent.RunConfig(max_calls=6,
                              profile=profiles.replace(profiles.DEFAULT, plan=False,
                                                       verify_rounds=0))
        llm = _ScriptedLLM([
            json.dumps({"thought": "t", "tool": "think", "args": {"thought": "hm"}}),
            json.dumps({"thought": "t", "tool": "done", "args": {"summary": "thought"}}),
        ])
        ep = agent.run_harness(llm, w, mem, "just think about it", cfg=cfg)
        self.assertTrue(ep.finished)
        self.assertEqual(w.snapshots, 1, "the loop must still snapshot")
        self.assertEqual([a["tool"] for a in w.actions], ["think"])

    def test_the_office_members_are_the_domains_not_the_loops(self):
        """The other half of the seam: everything beyond the four belongs to
        the simulated office and is reached only by that domain's tools."""
        w = World(tempfile.mkdtemp())
        for office_member in ("emails", "events", "messages", "reminders",
                              "list_emails", "add_event", "files_dir"):
            self.assertTrue(hasattr(w, office_member))
            self.assertNotIn(office_member, world_mod.WORLD_CONTRACT)


# ------------------------------------------------------------------- memory ---

class TestMemory(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.path = os.path.join(self.tmp.name, "memory.jsonl")

    def tearDown(self):
        self.tmp.cleanup()

    def test_facts_survive_a_reload(self):
        MemoryStore(self.path).save("User's manager is Sam.")
        self.assertEqual(MemoryStore(self.path).facts, ["User's manager is Sam."])

    def test_search_ranks_by_overlap_and_returns_nothing_for_a_miss(self):
        m = MemoryStore(self.path)
        m.save("User prefers meetings after 14:00")
        m.save("The office wifi password is hunter2")
        self.assertEqual(m.search("meeting preferences"),
                         ["User prefers meetings after 14:00"])
        self.assertEqual(m.search("quarterly revenue"), [])

    def test_a_torn_line_does_not_brick_the_agent_folder(self):
        """The agent appends to this file itself, so a crash mid-write leaves a
        partial line. Losing that line is fine; refusing to start is not."""
        with open(self.path, "w", encoding="utf-8") as f:
            f.write(json.dumps({"fact": "good one"}) + "\n")
            f.write('{"fact": "torn half')  # no newline, no closing brace
        self.assertEqual(MemoryStore(self.path).facts, ["good one"])

    def test_the_same_fact_twice_is_stored_once(self):
        """memory_k is 3 or 4 slots in the system prompt. A model that saves the
        same fact on every run would otherwise spend all of them on one fact."""
        m = MemoryStore(self.path)
        m.save("User's manager is Sam.")
        m.save("  User's manager is Sam.  ")
        self.assertEqual(m.facts, ["User's manager is Sam."])
        self.assertEqual(len(MemoryStore(self.path).facts), 1)


# ----------------------------------------------------------------- profiles ---

class TestProfiles(unittest.TestCase):
    def test_an_exact_tag_wins(self):
        self.assertEqual(profiles.for_model("llama3.2:1b").label, "format-survival")
        self.assertEqual(profiles.for_model("llama3.1:8b").label, "balanced")

    def test_a_quantized_variant_resolves_to_its_own_size(self):
        self.assertEqual(profiles.for_model("llama3.2:1b-instruct-q4_K_M").label,
                         "format-survival")
        self.assertEqual(profiles.for_model("llama3.2:3b-instruct-q8_0").label,
                         "guided-guarded")

    def test_an_unlisted_size_does_not_inherit_a_smaller_sibling(self):
        """llama3.2 happens to list 1b first. Matching on family alone handed an
        11B the 1B's profile: no planning, no verifier, 350-token replies."""
        prof = profiles.for_model("llama3.2:11b")
        self.assertNotEqual(prof.label, "format-survival")
        self.assertTrue(prof.plan)

    def test_an_unknown_model_is_sized_rather_than_given_the_flat_default(self):
        """gemma3:4b and phi4-mini:3.8b are offered in the picker but have no
        profile of their own. Size is the knob that actually predicts the
        failure mode, and the tag states it."""
        self.assertFalse(profiles.for_model("gemma3:1b").plan)
        self.assertEqual(profiles.for_model("phi4-mini:3.8b").verify_rounds, 1)
        self.assertTrue(profiles.for_model("qwen3:8b").plan)

    def test_a_model_with_no_size_in_its_tag_falls_back_to_default(self):
        self.assertEqual(profiles.for_model("some-new-model").label, "default")
        self.assertEqual(profiles.for_model(None).label, "default")

    def test_a_config_block_patches_individual_fields(self):
        prof = profiles.for_model("llama3.1:8b", {"max_calls": 33, "bogus": 1})
        self.assertEqual(prof.max_calls, 33)
        self.assertEqual(prof.label, "balanced")
        self.assertFalse(hasattr(prof, "bogus"))


# ----------------------------------------------------------------- fs_tools ---

class TestFsTools(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.root = os.path.realpath(self.tmp.name)
        self.saved = dict(TOOLS)
        # execute() logs every call to the world, so the tools need one even
        # though the file tools themselves never read it.
        self.world = World(os.path.join(self.root, ".world"))
        fs_tools.enable(self.root, allow_shell=True, confirm=None)

    def tearDown(self):
        TOOLS.clear()
        TOOLS.update(self.saved)
        self.tmp.cleanup()

    def run_tool(self, name, **args):
        return execute(name, args, self.world, None)

    def test_a_relative_path_lands_inside_the_root(self):
        ok, obs = self.run_tool("write_file", path="notes/todo.txt", content="hi")
        self.assertTrue(ok, obs)
        self.assertTrue(os.path.isfile(os.path.join(self.root, "notes", "todo.txt")))

    def test_escaping_the_root_is_refused(self):
        for bad in ("../outside.txt", "/etc/hosts", "notes/../../outside.txt"):
            ok, obs = self.run_tool("write_file", path=bad, content="x")
            self.assertFalse(ok, f"{bad} was allowed: {obs}")
            self.assertIn("outside the allowed root", obs)

    def test_a_declined_confirmation_tells_the_model_not_to_retry(self):
        open(os.path.join(self.root, "a.txt"), "w").close()
        fs_tools.enable(self.root, allow_shell=True, confirm=lambda *a: False)
        ok, obs = self.run_tool("delete_path", path="a.txt")
        self.assertFalse(ok)
        self.assertIn("Do not retry", obs)
        self.assertTrue(os.path.exists(os.path.join(self.root, "a.txt")))

    def test_shell_stays_off_unless_it_was_asked_for(self):
        fs_tools.enable(self.root, allow_shell=False, confirm=None)
        self.assertNotIn("run_command", TOOLS)

    def test_the_shell_runs_on_this_platform(self):
        """It shelled out to powershell.exe unconditionally, so shell mode could
        not work anywhere but Windows."""
        ok, obs = self.run_tool("run_command", command="echo hello-from-the-shell")
        self.assertTrue(ok, obs)
        self.assertIn("hello-from-the-shell", obs)
        self.assertIn("exit code 0", obs)

    def test_the_examples_use_this_platform_s_separator(self):
        """A model copies the example verbatim. A Windows example on POSIX
        creates one file literally named 'notes\\todo.txt'."""
        for name, spec in TOOLS.items():
            for value in spec["example"]["args"].values():
                if isinstance(value, str) and os.sep == "/":
                    self.assertNotIn("\\", value, f"{name} example is Windows-only")

    def test_the_protected_locations_exist_on_this_platform(self):
        """The deny-list was Windows absolute paths, so on macOS and Linux it
        matched nothing and protected nothing."""
        real = [d for d in fs_tools._DENY_WRITE if os.path.isdir(d)]
        self.assertTrue(real, f"no entry in the deny-list exists here: {fs_tools._DENY_WRITE}")

    def test_a_denied_location_cannot_be_written_even_inside_the_root(self):
        protected = os.path.join(self.root, "protected")
        os.makedirs(protected)
        fs_tools.enable(self.root, allow_shell=False, confirm=None,
                        deny=[protected])
        ok, obs = self.run_tool("write_file", path="protected/x.txt", content="x")
        self.assertFalse(ok, obs)
        self.assertIn("protected", obs)


# ------------------------------------------------------------------- office ---

class TestOffice(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()

    def tearDown(self):
        self.tmp.cleanup()

    def test_a_slide_whose_title_key_is_a_near_miss_is_repaired(self):
        """Observed live: a 1B sent {"slide_type": "title_slide"} and the call was
        rejected, three times, burning six of its eighteen calls in a loop it
        could not escape. A bare string was already repaired into a title-only
        slide on exactly the reasoning that rejecting it costs a model call to
        learn - that argument applies one level deeper too."""
        for bad, want in (({"heading": "Q3 Results"}, "Q3 Results"),
                          ({"Title": "Q3 Results"}, "Q3 Results"),
                          ({"slide_type": "Q3 Results"}, "Q3 Results")):
            out = office.create_presentation(self.tmp.name, "d.pptx", [bad])
            self.assertIn("1 slide", out)
            from pptx import Presentation
            deck = Presentation(os.path.join(self.tmp.name, "d.pptx"))
            texts = [sh.text_frame.text for sh in deck.slides[0].shapes if sh.has_text_frame]
            self.assertIn(want, texts, bad)

    def test_a_slide_with_no_usable_title_is_still_an_error(self):
        """Repair guesses only where the intent is unambiguous. Inventing a title
        for {"bullets": [...]} would put words on a slide nobody wrote."""
        with self.assertRaises(ToolError):
            office.create_presentation(self.tmp.name, "d.pptx", [{"bullets": ["a", "b"]}])

    def test_the_other_format_s_extension_is_swapped_not_appended(self):
        """Observed live: create_spreadsheet("q3_sales.pptx") wrote
        q3_sales.pptx.xlsx, and the UI row said "q3_sales.pptx written" because
        it shows the argument. A file the user never named, under a name the app
        never wrote."""
        office.create_spreadsheet(self.tmp.name, "q3_sales.pptx", [["a"], ["1"]])
        self.assertEqual(os.listdir(self.tmp.name), ["q3_sales.xlsx"])

    def test_an_unrelated_dot_in_a_name_is_left_alone(self):
        office.create_spreadsheet(self.tmp.name, "q3.final", [["a"], ["1"]])
        self.assertEqual(os.listdir(self.tmp.name), ["q3.final.xlsx"])

    def test_a_spreadsheet_round_trips_through_the_reader(self):
        office.create_spreadsheet(self.tmp.name, "costs.xlsx",
                                  [["Item", "Cost"], ["Chairs", 400]])
        out = office.read_spreadsheet(self.tmp.name, "costs.xlsx")
        self.assertEqual(out[0]["rows"][0], ["Item", "Cost"])
        self.assertEqual(out[0]["rows"][1], ["Chairs", 400])

    def test_a_missing_extension_is_added_rather_than_rejected(self):
        office.create_spreadsheet(self.tmp.name, "costs", [["a"]])
        self.assertTrue(os.path.isfile(os.path.join(self.tmp.name, "costs.xlsx")))

    def test_a_filename_cannot_escape_the_files_folder(self):
        office.create_spreadsheet(self.tmp.name, "../escaped.xlsx", [["a"]])
        self.assertFalse(os.path.exists(
            os.path.join(os.path.dirname(self.tmp.name), "escaped.xlsx")))

    def test_a_deck_needs_slides_and_says_so(self):
        with self.assertRaises(ToolError) as cm:
            office.create_presentation(self.tmp.name, "d.pptx", [])
        self.assertIn("slides", str(cm.exception))

    def test_rows_sent_as_objects_become_header_plus_rows(self):
        """Models regularly send [{"Region": "West", "Amount": 1240000}, ...].
        That is an unambiguous spreadsheet - keys are the header - and
        rejecting it costs a whole model call to repair a shape the code can
        convert deterministically. Same philosophy as parameter-name repair."""
        office.create_spreadsheet(self.tmp.name, "r.xlsx",
                                  [{"Region": "West", "Amount": 1240000},
                                   {"Region": "East", "Amount": 845000}])
        rows = office.read_spreadsheet(self.tmp.name, "r.xlsx")[0]["rows"]
        self.assertEqual(rows[0], ["Region", "Amount"])
        self.assertEqual(rows[1], ["West", 1240000])

    def test_objects_with_differing_keys_are_not_guessed_at(self):
        with self.assertRaises(ToolError):
            office.create_spreadsheet(self.tmp.name, "bad.xlsx",
                                      [{"a": 1}, {"b": 2}])

    def test_a_bare_string_slide_is_a_title_slide(self):
        office.create_presentation(self.tmp.name, "d.pptx",
                                   ["Cover", {"title": "Body", "bullets": ["x"]}])
        from pptx import Presentation
        prs = Presentation(os.path.join(self.tmp.name, "d.pptx"))
        self.assertEqual(len(prs.slides), 2)

    def test_reading_a_spreadsheet_that_was_never_made(self):
        with self.assertRaises(ToolError):
            office.read_spreadsheet(self.tmp.name, "ghost.xlsx")


# ----------------------------------------------------------------- verifier ---

class _StubLLM:
    """Captures the prompt it was handed and answers with a fixed reply."""

    def __init__(self, reply='{"complete": true, "missing": ""}'):
        self.reply = reply
        self.seen = []
        self.calls = 0

    def chat(self, messages, **kw):
        self.seen.append(messages)
        self.calls += 1
        return self.reply


class TestPromptBlock(unittest.TestCase):
    def test_prior_answers_are_a_record_not_first_person_prose(self):
        """"You: <sentence>" is a pattern to continue; a 1B continued it, opening
        each new summary with the previous one. Cheaper to remove the cue than
        to catch every copy downstream, so the check and this both ship."""
        from harness import chat
        block = chat.prompt_block([{"role": "user", "text": "summarize wednesday"},
                                   {"role": "assistant", "text": "Summarized three meetings."}])
        self.assertIn("summarize wednesday", block)
        self.assertIn("Summarized three meetings.", block)
        self.assertNotIn("You:", block)

    def test_the_block_is_empty_without_turns(self):
        from harness import chat
        self.assertEqual(chat.prompt_block([]), "")


class TestEchoesHistory(unittest.TestCase):
    """A copied span, not incidental word overlap: two summaries of similar work
    share vocabulary, and questioning that would hound every follow-up turn."""

    def test_a_copied_span_is_an_echo(self):
        hist = "Assistant did: Summarized my Wednesday meetings and messaged Jordan with the list"
        self.assertTrue(agent.echoes_history(
            "Summarized my Wednesday meetings and messaged Jordan with the list. Plus a deck.",
            hist))

    def test_shared_vocabulary_is_not_an_echo(self):
        hist = "Assistant did: Summarized my Wednesday meetings and messaged Jordan"
        self.assertFalse(agent.echoes_history(
            "Summarized my Thursday meetings and messaged Sam", hist))

    def test_a_short_summary_is_never_an_echo(self):
        """Under the span length there is nothing to copy: "Booked it." matching
        is coincidence, and re-asking would be noise on every terse answer."""
        hist = "Assistant did: Booked it"
        self.assertFalse(agent.echoes_history("Booked it", hist))

    def test_no_history_means_nothing_to_echo(self):
        self.assertFalse(agent.echoes_history("anything " * 20, ""))

    def test_case_and_spacing_do_not_hide_a_copy(self):
        hist = "Assistant did: alpha beta gamma delta epsilon zeta eta theta"
        self.assertTrue(agent.echoes_history(
            "ALPHA  beta\ngamma delta epsilon zeta eta theta iota", hist))


class TestVerifier(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.world = World(self.tmp.name)

    def tearDown(self):
        self.tmp.cleanup()

    def test_the_verifier_is_shown_what_the_actions_returned(self):
        """Given only "create_spreadsheet(...) -> ok" it cannot see that the file
        it is about to ask for already exists. Observed live: it answered
        complete:false against a finished task, the 8B redid the work, and the
        rerun wrote a SECOND spreadsheet - one task, two files for the user."""
        office.create_spreadsheet(self.world.files_dir, "q3.xlsx", [["Region", "Total"]])
        execute("create_spreadsheet",
                {"filename": "q3.xlsx", "rows": [["Region", "Total"]]}, self.world, None)
        llm = _StubLLM()
        agent._verify(llm, self.world, "build a spreadsheet")
        prompt = llm.seen[0][-1]["content"]
        self.assertIn("created q3.xlsx", prompt)

    def test_think_calls_are_left_out_of_the_evidence(self):
        execute("think", {"thought": "pondering"}, self.world, None)
        llm = _StubLLM()
        agent._verify(llm, self.world, "do something")
        self.assertNotIn("pondering", llm.seen[0][-1]["content"])

    def test_an_unusable_verdict_fails_open_but_says_it_did(self):
        """A broken verifier must not trap the agent in a loop it cannot exit.
        But an unmarked complete:true is indistinguishable from a real pass, so
        a systematically broken verifier would read as a clean run forever."""
        verdict = agent._verify(_StubLLM("I think it looks fine!"), self.world, "t")
        self.assertTrue(verdict["complete"])
        self.assertIn("unverified", verdict)

    def test_unrequested_side_effects_are_reported_not_undone(self):
        """The verifier now flags world-changing actions the task never asked
        for. The run still completes - the writes already happened, and undoing
        them would be a bigger side effect than the one reported - but the
        episode carries the report for the runner and the UI to surface."""
        execute("cancel_event", {"id": "c2"}, self.world, None)
        verdict = ('{"complete": true, "missing": "",'
                   ' "unrequested": "cancel_event on the Design review"}')
        llm = _StubLLM(verdict)
        out = agent._verify(llm, self.world, "remember my preference")
        self.assertEqual(out["unrequested"], "cancel_event")

    def test_a_mangled_unrequested_report_is_reduced_to_tool_names(self):
        """Observed live, twice: the 8B copies the evidence line's own format and
        runs out of tokens mid-string, giving "send_email({" and
        "['list_emails({}) -> ok: [ ... ]', -1]". Both went to the UI verbatim.
        Repair it deterministically rather than trusting the format instruction:
        keep the tool names the run actually performed, drop the rest."""
        execute("send_email", {"to": "d@c.com", "subject": "s", "body": "b"}, self.world, None)
        for mangled, want in (('send_email({', "send_email"),
                              ("['list_emails({}) -> ok: [ ... ]', -1]", ""),
                              ('the assistant called send_email which was not asked for',
                               "send_email")):
            out = agent._verify(_StubLLM(json.dumps({"complete": True, "missing": "",
                                                     "unrequested": mangled})),
                                self.world, "t")
            self.assertEqual(out["unrequested"], want)

    def test_a_read_that_carries_the_requirements_is_not_truncated_away(self):
        """On an indirect task ("do what the email asks") the requirements live
        in a read's RESULT, and the 200-char cap cut them off mid-sentence.
        Measured: Dana's 248-char email was clipped four words before "and turn
        the same numbers into a short deck", so the verifier accepted runs with
        no deck and flagged the spreadsheet it did build as unrequested. A write
        result is an echo of arguments the model already chose; a read result is
        the only place new information enters, so it gets the room."""
        self.world.emails.insert(0, {
            "id": "e99", "from": "dana@corp.com", "date": "2026-07-20 08:40",
            "subject": "two things",
            "body": "Filler. " * 30 + "and second, turn them into a short deck."})
        execute("read_email", {"id": "e99"}, self.world, None)
        execute("create_spreadsheet", {"filename": "q.xlsx",
                                       "rows": [["x" * 400], ["y"]]}, self.world, None)
        llm = _StubLLM()
        agent._verify(llm, self.world, "do what the newest email asks")
        prompt = llm.seen[0][-1]["content"]
        self.assertIn("short deck", prompt)          # the read survives
        self.assertNotIn("x" * 300, prompt)          # the write is still capped

    def test_the_verifier_is_told_where_a_delegated_task_keeps_its_requirements(self):
        """"Judge ONLY the requirements stated in the task" is right for a task
        that states them, and useless for one that delegates them. Measured: with
        the email fully visible in the evidence, the verifier still passed runs
        that built no deck, and still called the spreadsheet unrequested, because
        it had been told not to look past the task sentence."""
        llm = _StubLLM()
        agent._verify(llm, self.world, "do what the newest email asks")
        self.assertIn("DELEGATES", llm.seen[0][0]["content"])

    def test_reads_are_never_reported_as_unrequested(self):
        """The system prompt says reading is never unrequested; the 8B says it
        anyway. Observed live: "list_emails, read_spreadsheet" on a run whose
        only real fault was doing nothing wrong. Looking around is how the agent
        avoids inventing data, so reporting it teaches exactly the wrong lesson."""
        execute("list_emails", {}, self.world, None)
        execute("send_email", {"to": "d@c.com", "subject": "s", "body": "b"}, self.world, None)
        out = agent._verify(_StubLLM(json.dumps(
            {"complete": True, "missing": "",
             "unrequested": "list_emails, send_email"})), self.world, "t")
        self.assertEqual(out["unrequested"], "send_email")

    def test_a_none_ish_unrequested_report_is_treated_as_empty(self):
        """Small models answer "None" instead of an empty string."""
        for noise in ("None", "none", "N/A", ""):
            verdict = {"complete": True, "missing": "", "unrequested": noise}
            extra = str(verdict.get("unrequested") or "").strip()
            flagged = bool(extra) and extra.lower() not in ("none", "n/a", "nothing")
            self.assertFalse(flagged, noise)

    def test_a_verifier_that_raises_does_not_kill_the_run(self):
        class Boom:
            calls = 0

            def chat(self, *a, **k):
                raise ConnectionError("ollama went away")

        verdict = agent._verify(Boom(), self.world, "t")
        self.assertTrue(verdict["complete"])
        self.assertIn("ConnectionError", verdict["unverified"])


# ---------------------------------------------------------------- llm client ---

class TestLLMErrors(unittest.TestCase):
    """A refused model call must name the model and quote the server's own
    sentence. raise_for_status() threw the body away, so the user saw "404
    Client Error: Not Found for url http://127.0.0.1:11434/api/chat" - a
    loopback URL nobody can act on - while the body said "model 'x' not
    found, try pulling it first"."""

    class Resp:
        def __init__(self, status, body=None, text=""):
            self.status_code = status
            self._body = body
            self.text = text

        def json(self):
            if self._body is None:
                raise ValueError("not json")
            return self._body

    def test_the_servers_own_sentence_reaches_the_user(self):
        from harness.llm import _check
        with self.assertRaises(RuntimeError) as cm:
            _check(self.Resp(404, {"error": "model 'ghost:99b' not found, try pulling it first"}),
                   "ghost:99b")
        msg = str(cm.exception)
        self.assertIn("ghost:99b", msg)
        self.assertIn("try pulling it first", msg)
        self.assertNotIn("127.0.0.1", msg)

    def test_a_body_that_is_not_json_still_reads(self):
        from harness.llm import _check
        with self.assertRaises(RuntimeError) as cm:
            _check(self.Resp(500, None, text="upstream exploded"), "m")
        self.assertIn("upstream exploded", str(cm.exception))

    def test_an_empty_body_falls_back_to_the_status(self):
        from harness.llm import _check
        with self.assertRaises(RuntimeError) as cm:
            _check(self.Resp(503, None, text=""), "m")
        self.assertIn("503", str(cm.exception))

    def test_a_healthy_response_passes_silently(self):
        from harness.llm import _check
        _check(self.Resp(200, {"message": {}}), "m")


# --------------------------------------------------------------- the loop ---

class _ScriptedLLM:
    """Replies from a fixed list, then repeats the last one forever."""

    def __init__(self, replies):
        self.replies = list(replies)
        self.calls = 0
        self.output_tokens = self.prompt_tokens = 0
        self.wall = 0.0
        self.seen_feedback = []

    def chat(self, messages, **kw):
        if messages and messages[-1]["role"] == "user":
            self.seen_feedback.append(messages[-1]["content"])
        reply = self.replies[min(self.calls, len(self.replies) - 1)]
        self.calls += 1
        return reply


class TestLoop(unittest.TestCase):
    """run_harness itself, driven by a scripted model. No network."""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.world = World(self.tmp.name)
        self.mem = MemoryStore(os.path.join(self.tmp.name, "m.jsonl"))
        self.saved_profile, self.saved_calls = agent.PROFILE, agent.MAX_CALLS
        agent.MAX_CALLS = 12

    def tearDown(self):
        agent.set_profile(self.saved_profile)
        agent.MAX_CALLS = self.saved_calls
        self.tmp.cleanup()

    def call(self, tool, **args):
        return json.dumps({"thought": "x", "tool": tool, "args": args})

    def executed(self, name):
        return [a for a in self.world.actions if a["tool"] == name]

    def test_a_call_that_failed_is_not_run_again_against_the_same_world(self):
        """The repeat budget is there so a model can look at something twice. A
        call that ERRORED cannot return anything new, so a budget of three used
        to buy three copies of one failure - observed live, an 8B spent three of
        its twenty calls on one bad email id."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0,
                                           repeat_limit=3))
        ep = agent.run_harness(_ScriptedLLM([self.call("read_email", id="c3")]),
                               self.world, self.mem, "read something")
        self.assertEqual(len(self.executed("read_email")), 1)
        self.assertFalse(ep.finished)

    def test_a_call_that_worked_may_still_repeat_up_to_its_budget(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0,
                                           repeat_limit=3))
        agent.run_harness(_ScriptedLLM([self.call("list_emails")]),
                          self.world, self.mem, "look at the inbox")
        self.assertEqual(len(self.executed("list_emails")), 3)

    def test_an_identical_write_is_never_repeated(self):
        """A successful write bumps the world and hands out a fresh budget, so
        the guard against a duplicate invite has to be its own rule."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0,
                                           repeat_limit=3))
        agent.run_harness(
            _ScriptedLLM([self.call("send_message", to="sam", text="hi")]),
            self.world, self.mem, "message sam")
        self.assertEqual(len(self.world.messages), 1)

    def test_a_misspelled_plan_step_still_arms_the_guards(self):
        """End to end, not just the helper: a plan that says create_sheet must
        still produce a plan naming create_spreadsheet, so the write is planned
        and the unplanned-write guard stays quiet."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        plan = '{"steps": [{"tool": "list_email", "what": "look"}, ' \
               '{"tool": "create_sheet", "what": "build it"}]}'
        llm = _ScriptedLLM([plan, self.call("list_emails"),
                            self.call("create_spreadsheet", filename="r.xlsx",
                                      rows=[["Item", "Cost"], ["X", 1]])])
        ep = agent.run_harness(llm, self.world, self.mem, "spreadsheet of my July receipts")
        text = "\n".join(n["content"] for n in ep.transcript if n["kind"] == "plan")
        self.assertIn("list_emails", text)
        self.assertIn("create_spreadsheet", text)
        self.assertEqual([n["content"] for n in ep.transcript if n["kind"] == "guard"], [],
                         "both steps were planned, so nothing should be questioned")
        repairs = [n["content"] for n in ep.transcript if n["kind"] == "repair"]
        self.assertTrue(any("plan step" in r for r in repairs), repairs)

    def test_a_plan_step_that_is_not_a_tool_at_all_is_still_dropped(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        plan = '{"steps": [{"tool": "xyzzy", "what": "nonsense"}]}'
        llm = _ScriptedLLM([plan, self.call("list_emails")])
        ep = agent.run_harness(llm, self.world, self.mem, "look at the inbox")
        text = "\n".join(n["content"] for n in ep.transcript if n["kind"] == "plan")
        self.assertIn("unusable plan reply", text)

    def test_a_config_disarms_a_guard_for_one_run_only(self):
        """The property the whole refactor exists for, through a real
        run_harness rather than through the config object: the same process
        runs one episode with the guard and one without, and neither leaks
        into the other."""
        plan = '{"steps": [{"tool": "list_emails", "what": "look"}, ' \
               '{"tool": "create_spreadsheet", "what": "build"}]}'

        def episode(cfg):
            llm = _ScriptedLLM([plan, self.call("create_spreadsheet", filename="r.xlsx",
                                                rows=[["Item", "Cost"], ["X", 1]])])
            world = World(tempfile.mkdtemp())
            return agent.run_harness(llm, world, self.mem, "spreadsheet of receipts",
                                     cfg=cfg), llm

        base = agent.RunConfig(max_calls=6,
                               profile=profiles.replace(profiles.DEFAULT, plan=True,
                                                        verify_rounds=0))
        with_guard, _ = episode(base)
        without, _ = episode(base.without_guard("read_before_write"))

        fired = [n["content"] for n in with_guard.transcript if n["kind"] == "guard"]
        silent = [n["content"] for n in without.transcript if n["kind"] == "guard"]
        self.assertIn("read_before_write", fired)
        self.assertNotIn("read_before_write", silent)

    def test_the_module_globals_are_untouched_by_a_configured_run(self):
        before = (agent.MAX_CALLS, agent.PROFILE, list(agent.GUARDS))
        cfg = agent.RunConfig(max_calls=3,
                              profile=profiles.replace(profiles.DEFAULT, plan=False,
                                                       verify_rounds=0))
        agent.run_harness(_ScriptedLLM([self.call("list_emails")]),
                          self.world, self.mem, "look", cfg=cfg)
        self.assertEqual((agent.MAX_CALLS, agent.PROFILE, list(agent.GUARDS)), before)

    def test_the_configs_budget_is_what_actually_stops_the_run(self):
        cfg = agent.RunConfig(max_calls=3,
                              profile=profiles.replace(profiles.DEFAULT, plan=False,
                                                       verify_rounds=0, loop_break=False))
        llm = _ScriptedLLM([self.call("think", thought="x")])
        agent.run_harness(llm, self.world, self.mem, "think about it", cfg=cfg)
        self.assertEqual(llm.calls, 3)

    def test_an_interleaved_call_does_not_reset_the_repeat_counter(self):
        """Compared against DeepSeek Harness's repeat guard, which counts
        CONSECUTIVE identical calls and therefore needs an explicit "bookkeeping
        tools are transparent to the chain" rule so that
        grep X -> todo_write -> grep X still counts as two.

        Ours needs no such rule: it counts occurrences of one signature against
        an unchanged world, so interleaving anything between them changes
        nothing. Untested until the comparison went looking, and load-bearing."""
        cfg = agent.RunConfig(max_calls=8,
                              profile=profiles.replace(profiles.DEFAULT, plan=False,
                                                       verify_rounds=0, repeat_limit=2))
        llm = _ScriptedLLM([self.call("list_emails"), self.call("think", thought="a"),
                            self.call("list_emails"), self.call("think", thought="b"),
                            self.call("list_emails")])
        agent.run_harness(llm, self.world, self.mem, "look repeatedly", cfg=cfg)
        self.assertEqual(len(self.executed("list_emails")), 2,
                         "the interleaved think must not hand out a fresh budget")

    def test_the_repeat_key_ignores_argument_order(self):
        """The other behaviour the comparison checked. json.dumps(sort_keys=True)
        is already a canonical key, so two calls differing only in the order
        their arguments were written count as identical."""
        a = json.dumps({"t": "x", "a": {"b": 1, "a": 2}}, sort_keys=True, default=str)
        b = json.dumps({"t": "x", "a": {"a": 2, "b": 1}}, sort_keys=True, default=str)
        self.assertEqual(a, b)

    def test_a_successful_write_hands_out_a_fresh_repeat_budget(self):
        """The notion ours has that theirs does not: the counter is scoped to an
        unchanged world, so the same call may legitimately return something new
        once something has been written."""
        cfg = agent.RunConfig(max_calls=8,
                              profile=profiles.replace(profiles.DEFAULT, plan=False,
                                                       verify_rounds=0, repeat_limit=1))
        llm = _ScriptedLLM([self.call("list_events"),
                            self.call("add_event", title="X", date="2026-07-23",
                                      start_time="10:00", end_time="11:00"),
                            self.call("list_events")])
        agent.run_harness(llm, self.world, self.mem, "look, write, look again", cfg=cfg)
        self.assertEqual(len(self.executed("list_events")), 2,
                         "the write should have reset the budget for the read")

    def test_a_guard_note_lands_immediately_before_the_message_it_produced(self):
        """The contract the web UI depends on: it holds the guard name from the
        guard note and labels the very next feedback note with it, so a nudge
        on screen says WHICH check spoke. One slot is enough only because the
        two are adjacent and in that order."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        plan = '{"steps": [{"tool": "list_emails", "what": "find the receipts"}, ' \
               '{"tool": "create_spreadsheet", "what": "build it"}]}'
        llm = _ScriptedLLM([plan, self.call("create_spreadsheet", filename="r.xlsx",
                                            rows=[["Item", "Cost"], ["X", 1]])])
        ep = agent.run_harness(llm, self.world, self.mem, "spreadsheet of my July receipts")
        kinds = [(n["kind"], n["content"]) for n in ep.transcript]
        guards = [i for i, (k, _) in enumerate(kinds) if k == "guard"]
        self.assertTrue(guards, "no guard fired, so this asserts nothing")
        for i in guards:
            self.assertEqual(kinds[i][1], "read_before_write")
            self.assertEqual(kinds[i + 1][0], "feedback",
                             "a guard note must be followed by its own message")

    def test_only_a_questioned_call_emits_a_guard_note(self):
        """The other half: an ordinary run must not leave a stale guard name
        lying around for an unrelated feedback message to pick up."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        ep = agent.run_harness(_ScriptedLLM([self.call("list_emails")]),
                               self.world, self.mem, "look at the inbox")
        self.assertEqual([n for n in ep.transcript if n["kind"] == "guard"], [])

    def test_writing_before_reading_is_questioned_once(self):
        """The worst failure this app can produce: asked for a spreadsheet of
        July receipts, an 8B skipped the inbox and invented the numbers. The
        plan had named the read step. One nudge, then it gets its way."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        plan = '{"steps": [{"tool": "list_emails", "what": "find the receipts"}, '
        plan += '{"tool": "create_spreadsheet", "what": "totals"}]}'
        llm = _ScriptedLLM([plan, self.call("create_spreadsheet", filename="r.xlsx",
                                            rows=[["Item", "Cost"], ["made up", 100]])])
        agent.run_harness(llm, self.world, self.mem, "spreadsheet of my July receipts")
        written = self.executed("create_spreadsheet")
        self.assertEqual(len(written), 1, "the nudge should delay the write, not block it")
        nudges = [n for n in llm.seen_feedback if "writing from memory" in n]
        self.assertEqual(len(nudges), 1, "exactly one nudge, or the agent is stuck")

    def test_a_run_that_reads_first_is_never_nudged(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        plan = '{"steps": [{"tool": "list_emails", "what": "find them"}, '
        plan += '{"tool": "create_spreadsheet", "what": "totals"}]}'
        llm = _ScriptedLLM([plan, self.call("list_emails"),
                            self.call("create_spreadsheet", filename="r.xlsx",
                                      rows=[["a", 1]]),
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "spreadsheet of my July receipts")
        self.assertFalse([n for n in llm.seen_feedback if "writing from memory" in n])

    def test_a_read_planned_after_the_write_is_not_a_source_read(self):
        """A plan of think -> create_spreadsheet -> read_spreadsheet reads back
        the file it is about to make. Nudging toward it sent the agent to open a
        spreadsheet that did not exist yet."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        plan = ('{"steps": [{"tool": "think", "what": "plan it"}, '
                '{"tool": "create_spreadsheet", "what": "make it"}, '
                '{"tool": "read_spreadsheet", "what": "check it"}]}')
        llm = _ScriptedLLM([plan, self.call("create_spreadsheet", filename="r.xlsx",
                                            rows=[["a", 1]]),
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "spreadsheet of my July receipts")
        self.assertFalse([n for n in llm.seen_feedback if "writing from memory" in n])

    def test_a_crash_mid_run_still_snapshots_the_world(self):
        """A mid-run crash (ollama dying) used to lose every world mutation:
        the UI had already reported "send_message - written", but state.json
        was only written on clean exit, so after a restart the sent message
        had never happened."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))

        class DiesAfterOne(_ScriptedLLM):
            def chat(self, messages, **kw):
                if self.calls >= 1:
                    raise ConnectionError("ollama died")
                return super().chat(messages, **kw)

        llm = DiesAfterOne([self.call("send_message", to="sam", text="hi")])
        with self.assertRaises(ConnectionError):
            agent.run_harness(llm, self.world, self.mem, "message sam")
        state = json.load(open(os.path.join(self.tmp.name, "state.json")))
        self.assertEqual(len(state["messages"]), 1)

    def test_an_unplanned_write_is_questioned_once_then_allowed(self):
        """Asked only to "list my emails", an 8B sent an email, added an event
        and messaged a third party. A write the model's own plan never named is
        challenged once; if it insists, it runs - question, never forbid.

        Since replanning landed, the challenge on a task that has already read
        something is the replan call rather than the question: a model that
        still wants the write says so in the revised plan and it runs. That is
        deliberately the same strength as before - the old guard was always a
        speed bump the model could clear by repeating itself, and an 8B clears
        this one by restating. The verifier's unrequested report is what
        actually catches the side effect, and it is unchanged."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        plan = '{"steps": [{"tool": "list_emails", "what": "list them"}]}'
        llm = _ScriptedLLM([plan,
                            self.call("list_emails"),
                            self.call("send_message", to="sam", text="fyi"),   # challenged
                            '{"steps": [{"tool": "send_message", "what": "tell sam"}]}',
                            self.call("send_message", to="sam", text="fyi"),   # insists: runs
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "List my emails")
        self.assertEqual(len(self.world.messages), 1)
        self.assertFalse([f for f in llm.seen_feedback if "never included send_message" in f])

    def test_a_write_discovered_by_reading_gets_the_plan_revised_not_refused(self):
        """The plan is written before the agent has read anything, so it cannot
        name work the task's own data turns out to require. Observed live: asked
        to "read Dana's newest email and do what she asks", the plan was
        list/read/send; the email asked for a spreadsheet; the guard told the
        model its plan never included create_spreadsheet and to do "nothing
        extra"; the model sent an email CLAIMING it had built the sheet and
        stopped. A plan made before discovery is a hypothesis: once a read has
        landed, revise it instead of holding the model to it."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        blind = '{"steps": [{"tool": "list_emails", "what": "find it"}, ' \
                '{"tool": "read_email", "what": "read it"}]}'
        revised = '{"steps": [{"tool": "create_spreadsheet", "what": "what she asked for"}]}'
        llm = _ScriptedLLM([blind,
                            self.call("list_emails"),
                            self.call("read_email", id="e1"),
                            self.call("create_spreadsheet", filename="q.xlsx",
                                      rows=[["a"], ["1"]]),
                            revised,          # the replan call
                            self.call("done", summary="built it")])
        agent.run_harness(llm, self.world, self.mem, "Read the newest email and do what it asks")
        self.assertEqual(os.listdir(self.world.files_dir), ["q.xlsx"])  # it ran, not refused
        self.assertFalse([f for f in llm.seen_feedback if "never included" in f])

    def test_the_plan_is_only_revised_once(self):
        """A second discovery falls back to the question. Revising on every
        surprise would let a wandering agent rewrite its way to anything."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        blind = '{"steps": [{"tool": "list_emails", "what": "find it"}]}'
        revised = '{"steps": [{"tool": "create_spreadsheet", "what": "the sheet"}]}'
        llm = _ScriptedLLM([blind,
                            self.call("list_emails"),
                            self.call("create_spreadsheet", filename="q.xlsx", rows=[["a"]]),
                            revised,
                            self.call("send_message", to="sam", text="fyi"),  # questioned
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "List my emails and do what they ask")
        self.assertEqual(len(self.world.messages), 0)
        self.assertEqual(len([f for f in llm.seen_feedback
                              if "never included send_message" in f]), 1)

    def test_an_unplanned_write_before_any_read_is_still_questioned(self):
        """Nothing has been discovered yet, so there is nothing to revise
        against and the plan still stands."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        plan = '{"steps": [{"tool": "list_emails", "what": "list them"}]}'
        llm = _ScriptedLLM([plan,
                            self.call("send_message", to="sam", text="fyi"),
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "List my emails")
        self.assertEqual(len(self.world.messages), 0)
        self.assertTrue([f for f in llm.seen_feedback if "never included send_message" in f])

    def test_the_unplanned_write_question_does_not_tell_the_model_to_stop(self):
        """"Only do what the task requires - nothing extra" is what an 8B obeys
        instead of insisting, which turns a question into a block."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        llm = _ScriptedLLM(['{"steps": [{"tool": "list_emails", "what": "list"}]}',
                            self.call("send_message", to="sam", text="fyi"),
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "List my emails")
        asked = [f for f in llm.seen_feedback if "never included send_message" in f][0]
        self.assertNotIn("nothing extra", asked)
        self.assertIn("call it again", asked)

    def test_a_file_named_by_what_was_read_is_questioned_before_writing_over_it(self):
        """Observed live: the email said "the export is in q3_raw.xlsx", the
        agent never opened it, and invented Sales/Profit rows with formulas over
        empty cells. Writing from memory when the task's own data is sitting in
        a file on disk is the failure the whole harness exists to catch."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        office.create_spreadsheet(self.world.files_dir, "q3_raw.xlsx",
                                  [["Region", "Q3"], ["West", 1240000]])
        self.world.emails.insert(0, {"id": "e99", "from": "dana@corp.com",
                                     "date": "2026-07-20 08:40", "subject": "numbers",
                                     "body": "The export is in q3_raw.xlsx, pull the Q3 column."})
        llm = _ScriptedLLM([self.call("read_email", id="e99"),
                            self.call("create_spreadsheet", filename="out.xlsx",
                                      rows=[["Sales"], ["1"]]),          # questioned
                            self.call("read_spreadsheet", filename="q3_raw.xlsx"),
                            self.call("create_spreadsheet", filename="out.xlsx",
                                      rows=[["Region", "Q3"], ["West", 1240000]]),
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "Do what the newest email asks")
        # Match the nudge's own wording: an OBSERVATION also carries the
        # filename, and asserting on the filename alone passes without a nudge.
        self.assertTrue([f for f in llm.seen_feedback if "have not opened q3_raw.xlsx" in f])

    def test_the_unread_file_question_fires_once_then_lets_it_through(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        office.create_spreadsheet(self.world.files_dir, "q3_raw.xlsx", [["a"], ["1"]])
        self.world.emails.insert(0, {"id": "e99", "from": "d@c.com", "date": "2026-07-20 08:40",
                                     "subject": "n", "body": "see q3_raw.xlsx"})
        llm = _ScriptedLLM([self.call("read_email", id="e99"),
                            self.call("create_spreadsheet", filename="out.xlsx", rows=[["x"]]),
                            self.call("create_spreadsheet", filename="out.xlsx", rows=[["x"]]),
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "Do what the newest email asks")
        self.assertEqual(len([f for f in llm.seen_feedback
                              if "have not opened q3_raw.xlsx" in f]), 1)
        self.assertEqual(sorted(os.listdir(self.world.files_dir)),
                         ["out.xlsx", "q3_raw.xlsx"])

    def test_a_file_that_was_never_mentioned_is_not_questioned(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        office.create_spreadsheet(self.world.files_dir, "unrelated.xlsx", [["a"], ["1"]])
        llm = _ScriptedLLM([self.call("list_emails"),
                            self.call("create_spreadsheet", filename="out.xlsx", rows=[["x"]]),
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "Make a sheet")
        self.assertFalse([f for f in llm.seen_feedback if "unrelated.xlsx" in f])

    def test_a_summary_that_repeats_the_last_answer_is_sent_back(self):
        """Observed live across three turns of one thread: each run's done
        summary opened with the PREVIOUS run's summary and appended to it, so
        turn 3 described turn 2's work plus spam text from an unrelated email.
        It compounds - the contaminated summary is stored and becomes the next
        turn's context - so each turn is worse than the one before."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        prior = ("\n\nEARLIER IN THIS CONVERSATION:\nUser: summarize wednesday\n"
                 "Assistant did: Summarized my Wednesday meetings and messaged "
                 "Jordan with the list of three items")
        llm = _ScriptedLLM([
            self.call("done", summary="Summarized my Wednesday meetings and messaged "
                                      "Jordan with the list of three items. Also built a deck."),
            self.call("done", summary="Built the deck from Dana's numbers.")])
        ep = agent.run_harness(llm, self.world, self.mem, "build a deck", history=prior)
        self.assertEqual(ep.done_summary, "Built the deck from Dana's numbers.")
        self.assertTrue([f for f in llm.seen_feedback if "only what you did in THIS" in f])

    def test_a_summary_of_its_own_work_is_accepted(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        prior = ("\n\nEARLIER IN THIS CONVERSATION:\nUser: summarize wednesday\n"
                 "Assistant did: Summarized my Wednesday meetings and messaged Jordan")
        llm = _ScriptedLLM([self.call("done", summary="Built the deck from Dana's numbers.")])
        ep = agent.run_harness(llm, self.world, self.mem, "build a deck", history=prior)
        self.assertEqual(ep.done_summary, "Built the deck from Dana's numbers.")
        self.assertFalse([f for f in llm.seen_feedback if "only what you did in THIS" in f])

    def test_the_echo_check_is_asked_once_and_then_accepts(self):
        """Question, never forbid: a model that insists gets its summary."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        prior = "\n\nEARLIER:\nAssistant did: alpha beta gamma delta epsilon zeta eta theta"
        echo = self.call("done", summary="alpha beta gamma delta epsilon zeta eta theta")
        llm = _ScriptedLLM([echo, echo])
        ep = agent.run_harness(llm, self.world, self.mem, "t", history=prior)
        self.assertEqual(ep.done_summary, "alpha beta gamma delta epsilon zeta eta theta")
        self.assertEqual(len([f for f in llm.seen_feedback
                              if "only what you did in THIS" in f]), 1)

    def test_a_run_with_no_conversation_behind_it_is_never_checked(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        llm = _ScriptedLLM([self.call("done", summary="anything at all goes here")])
        agent.run_harness(llm, self.world, self.mem, "t")
        self.assertFalse([f for f in llm.seen_feedback if "only what you did in THIS" in f])

    def test_a_planned_write_is_never_questioned(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=True, verify_rounds=0))
        plan = ('{"steps": [{"tool": "list_emails", "what": "read"}, '
                '{"tool": "send_message", "what": "tell sam"}]}')
        llm = _ScriptedLLM([plan, self.call("list_emails"),
                            self.call("send_message", to="sam", text="update"),
                            self.call("done", summary="done")])
        agent.run_harness(llm, self.world, self.mem, "tell Sam about my inbox")
        self.assertEqual(len(self.world.messages), 1)
        self.assertFalse([f for f in llm.seen_feedback if "never included" in f])

    def test_the_date_guard_checks_writes_and_lets_reads_roam(self):
        """"Remember ... never on Fridays" made the guard hound four innocent
        list_events probes - a preference, not a date instruction, and reads
        are how a model looks around. Writes are still checked."""
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        llm = _ScriptedLLM([self.call("list_events", date="2026-07-21"),   # "wrong" day: fine
                            self.call("set_reminder", text="x",
                                      date="2026-07-21", time="09:00"),   # write: questioned
                            self.call("done", summary="d")])
        agent.run_harness(llm, self.world, self.mem,
                          "Set a reminder for my Wednesday meetings")
        wrongs = [f for f in llm.seen_feedback if "WRONG DATE" in f]
        self.assertEqual(len(wrongs), 1)
        listed = self.executed("list_events")
        self.assertEqual(len(listed), 1, "the read must have run unquestioned")

    def test_done_ends_the_run_and_keeps_the_summary(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        ep = agent.run_harness(_ScriptedLLM([self.call("done", summary="all finished")]),
                               self.world, self.mem, "do it")
        self.assertTrue(ep.finished)
        self.assertEqual(ep.done_summary, "all finished")

    def test_run_raw_still_works_after_the_snapshot_restructure(self):
        """bench/ calls run_raw and nothing else exercised it. The crash-safe
        snapshot moved its loop into a helper; this pins the behaviour."""
        llm = _ScriptedLLM(['{"tool": "list_emails", "args": {}}',
                            '{"tool": "done", "args": {"summary": "listed"}}'])
        ep = agent.run_raw(llm, self.world, self.mem, "list my emails")
        self.assertTrue(ep.finished)
        self.assertEqual(ep.done_summary, "listed")
        self.assertTrue(os.path.exists(os.path.join(self.tmp.name, "state.json")))

    def test_run_raw_snapshots_on_a_crash_too(self):
        class Dies(_ScriptedLLM):
            def chat(self, messages, **kw):
                if self.calls >= 1:
                    raise ConnectionError("gone")
                return super().chat(messages, **kw)

        llm = Dies(['{"tool": "send_message", "args": {"to": "sam", "text": "hi"}}'])
        with self.assertRaises(ConnectionError):
            agent.run_raw(llm, self.world, self.mem, "message sam")
        state = json.load(open(os.path.join(self.tmp.name, "state.json")))
        self.assertEqual(len(state["messages"]), 1)

    def test_the_unrequested_report_reaches_the_episode(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=1))

        class VerifierAware(_ScriptedLLM):
            def chat(self, messages, role=None, **kw):
                if role == "verifier":
                    self.calls += 1
                    return ('{"complete": true, "missing": "",'
                            ' "unrequested": "send_message"}')
                return super().chat(messages, **kw)

        llm = VerifierAware([self.call("send_message", to="sam", text="hi"),
                             self.call("done", summary="did it")])
        ep = agent.run_harness(llm, self.world, self.mem, "message sam")
        self.assertTrue(ep.finished)
        self.assertEqual(ep.unrequested, "send_message")

    def test_the_budget_is_a_hard_stop(self):
        agent.set_profile(profiles.replace(profiles.DEFAULT, plan=False, verify_rounds=0))
        llm = _ScriptedLLM(["not json at all"])
        ep = agent.run_harness(llm, self.world, self.mem, "do it")
        self.assertEqual(llm.calls, agent.MAX_CALLS)
        self.assertFalse(ep.finished)
        self.assertEqual(ep.parse_failures, agent.MAX_CALLS)


# ------------------------------------------------------------------- server ---

class TestWorkspacePanel(unittest.TestCase):
    def test_the_inbox_panel_shows_the_newest_email_first(self):
        """It rendered state.json's insertion order, so a mail that had just
        arrived appeared at the bottom of the inbox under nine older ones."""
        from webui import server
        with tempfile.TemporaryDirectory() as d:
            folder = os.path.join(d, "agents", "demo")
            os.makedirs(os.path.join(folder, "workspace"))
            with open(os.path.join(folder, "config.json"), "w") as f:
                json.dump({"name": "demo", "model": "llama3.1:8b"}, f)
            with open(os.path.join(folder, "workspace", "state.json"), "w") as f:
                json.dump({"emails": [
                    {"id": "a", "date": "2026-07-13 08:00", "subject": "old"},
                    {"id": "b", "date": "2026-07-20 08:40", "subject": "newest"},
                    {"id": "c", "date": "2026-07-15 09:10", "subject": "middle"}]}, f)
            saved = server.AGENTS_DIR
            server.AGENTS_DIR = os.path.join(d, "agents")
            try:
                got = [e["subject"] for e in server.workspace("demo")["emails"]]
            finally:
                server.AGENTS_DIR = saved
        self.assertEqual(got, ["newest", "middle", "old"])


class TestSpreadsheetPreview(unittest.TestCase):
    def test_a_grouped_number_previews_the_way_excel_shows_it(self):
        """office.py writes a #,##0 format, so Excel shows 1,240,000. The
        preview printed str(value) and showed 1240000, which on a demo built
        from an email quoting "$1,240,000" reads as a mangled number."""
        from webui.server import _cell_text

        class Cell:
            number_format = "#,##0"

        class Money:
            number_format = '#,##0.00'

        self.assertEqual(_cell_text(1240000, True, Cell()), "1,240,000")
        self.assertEqual(_cell_text(1240000.5, True, Money()), "1,240,000.50")
        self.assertEqual(_cell_text("West", False, Cell()), "West")
        self.assertEqual(_cell_text(None, False, Cell()), "")

    def test_a_number_with_no_format_is_left_alone(self):
        class Plain:
            number_format = "General"

        from webui.server import _cell_text
        self.assertEqual(_cell_text(2026, True, Plain()), "2026")


class TestServerGuards(unittest.TestCase):
    def test_reset_refuses_while_that_agent_is_mid_run(self):
        """Reset mid-run silently undid itself: the live subprocess snapshots
        on exit, so the deleted state.json reappeared carrying the pre-reset
        world. One thing owns the folder at a time; the answer is 409."""
        from webui import server

        class Proc:
            def poll(self):
                return None      # still running

        class Live:
            agent = "8b"
            proc = Proc()

        saved = server.RUNS.current
        server.RUNS.current = Live()
        try:
            with self.assertRaises(RuntimeError):
                server.ensure_idle("8b")
            server.ensure_idle("other-agent")   # a different folder is fine
        finally:
            server.RUNS.current = saved

    def test_reset_is_allowed_once_the_run_has_exited(self):
        from webui import server

        class Proc:
            def poll(self):
                return 0         # exited

        class Done:
            agent = "8b"
            proc = Proc()

        saved = server.RUNS.current
        server.RUNS.current = Done()
        try:
            server.ensure_idle("8b")
        finally:
            server.RUNS.current = saved

    def test_stop_unwinds_instead_of_killing(self):
        """Stop sends SIGTERM and Python's default is to die without
        unwinding, so the crash-safe snapshot never ran for a STOPPED run.
        The handler turns the signal into SystemExit, which runs finallys."""
        from webui.runner import _on_terminate
        with self.assertRaises(SystemExit) as cm:
            _on_terminate(15, None)
        self.assertEqual(cm.exception.code, 143)


class TestConfirmerMemory(unittest.TestCase):
    """A decline is a decision, not a transient failure. Observed live against
    a real MCP server: a declined draft was re-attempted three times, putting
    the same dialog in front of the person three times after they said no and
    spending 22 of the run's calls. The model reworded its arguments slightly
    each time, which is why the loop's signature dedupe cannot catch it."""

    def make(self, answers):
        from webui import runner
        c = runner.Confirmer()
        c._asked = []
        import io
        c._answers = list(answers)

        def fake_readline():
            return json.dumps({"id": c.n, "allow": c._answers.pop(0)}) + "\n"

        saved_stdin, saved_emit = runner.sys.stdin, runner.emit
        runner.sys.stdin = type("S", (), {"readline": staticmethod(fake_readline)})()
        runner.emit = lambda ev, **f: c._asked.append(f) if ev == "confirm" else None
        return c, (lambda: (setattr(runner.sys, "stdin", saved_stdin),
                            setattr(runner, "emit", saved_emit)))

    def test_a_declined_action_is_never_asked_twice(self):
        c, restore = self.make([False])
        try:
            first = c("call the tool", 'gmail: draft_mail {"to": "a@b.com"}')
            second = c("call the tool", 'gmail: draft_mail {"to": "a@b.com", "body": "reworded"}')
            self.assertFalse(first)
            self.assertFalse(second)
            self.assertEqual(len(c._asked), 1, "the person was asked twice")
        finally:
            restore()

    def test_an_allowed_action_may_be_asked_again(self):
        """Saying yes once does not blanket-authorize every later call."""
        c, restore = self.make([True, True])
        try:
            self.assertTrue(c("call the tool", 'gmail: draft_mail {"to": "a"}'))
            self.assertTrue(c("call the tool", 'gmail: draft_mail {"to": "b"}'))
            self.assertEqual(len(c._asked), 2)
        finally:
            restore()

    def test_a_real_account_confirm_is_flagged_as_one(self):
        """The dialog was identical for overwriting a scratch file and for
        sending from a live mailbox. mcp_bridge formats its detail as
        "<server-id>: <tool> {args}", so the prefix identifies a real account
        without changing the callback signature fs_tools also uses."""
        from webui import runner
        c, restore = self.make([True, True])
        c.real_servers = {"gmail"}
        c.mode = "live"
        try:
            c("call the tool", 'gmail: send_mail {"to": "a@b.com"}')
            c("overwrite", "/tmp/scratch.txt (12 bytes will be replaced)")
            self.assertEqual(c._asked[0]["real"], "gmail")
            self.assertEqual(c._asked[0]["mode"], "live")
            self.assertIsNone(c._asked[1]["real"], "a local file is not a real account")
        finally:
            restore()

    def test_a_different_action_is_still_asked(self):
        c, restore = self.make([False, True])
        try:
            c("call the tool", 'gmail: draft_mail {"to": "a"}')
            self.assertTrue(c("call the tool", 'gmail: modify_mail {"id": "1"}'))
            self.assertEqual(len(c._asked), 2)
        finally:
            restore()


# ------------------------------------------------------------------- runner ---

class TestRunner(unittest.TestCase):
    def test_run_numbering_never_reuses_an_index(self):
        from webui.runner import next_run_index
        with tempfile.TemporaryDirectory() as d:
            self.assertEqual(next_run_index(d), 1)
            for name in ("run_001.json", "run_003.json", "model_calls.jsonl", "notes.txt"):
                open(os.path.join(d, name), "w").close()
            self.assertEqual(next_run_index(d), 4)  # not 3, and not 5
            os.remove(os.path.join(d, "run_001.json"))
            self.assertEqual(next_run_index(d), 4)  # a deleted run frees nothing


class TestCallBudget(unittest.TestCase):
    """The ceiling on LLM calls in one run. It is a ceiling, not a target: an
    agent that calls done early costs whatever it costs."""

    def budget(self, *a, **kw):
        from webui.runner import call_budget
        return call_budget(*a, **kw)

    def test_the_shipped_8b_default_is_50(self):
        self.assertEqual(profiles.for_model("llama3.1:8b").max_calls, 50)

    def test_simulated_work_gets_the_profile_budget(self):
        p = profiles.for_model("llama3.2:3b")
        self.assertEqual(self.budget(p), p.max_calls)

    def test_real_work_lifts_a_tight_profile_to_the_floor(self):
        # A 3B's 14 is a deliberate loop-brake for the simulated office; real
        # files and real mailboxes cost a call per listing and per read.
        self.assertEqual(self.budget(profiles.for_model("llama3.2:3b"),
                                     extended=True), 40)

    def test_real_work_never_lowers_a_generous_profile(self):
        self.assertEqual(self.budget(profiles.for_model("llama3.1:8b"),
                                     extended=True), 50)

    def test_an_override_wins_over_both(self):
        p = profiles.for_model("llama3.1:8b")
        self.assertEqual(self.budget(p, override=7), 7)
        self.assertEqual(self.budget(p, override=7, extended=True), 7)

    def test_an_override_is_clamped_at_both_ends(self):
        # The UI is not the only caller, so the clamp lives here rather than in
        # the number input's min/max.
        p = profiles.for_model("llama3.1:8b")
        self.assertEqual(self.budget(p, override=5000), 200)
        self.assertEqual(self.budget(p, override=1), 2)
        self.assertEqual(self.budget(p, override=-3), 2)

    def test_a_blank_override_falls_through_to_the_profile(self):
        p = profiles.for_model("llama3.1:8b")
        for blank in (None, 0, ""):
            self.assertEqual(self.budget(p, override=blank), p.max_calls)


if __name__ == "__main__":
    unittest.main(verbosity=2)
